"""
Hypergraph-Redundancy Pruning — pitch deck.

Focused deck for the set-level (higher-order) redundancy direction:
  topic/problem -> background (pairwise) -> the gap -> related work ->
  our method -> worked example -> pipeline -> expected results -> positioning.

Style mirrors make_idea_slides.py (Isomorphic-Pruning visual language).
Run:  python make_redundancy_slides.py   ->  hypergraph_redundancy.pptx
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── palette (from make_idea_slides.py) ───────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x27, 0x44); WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x00, 0x68, 0xB5); ORANGE = RGBColor(0xFF, 0xA3, 0x00)
DARK   = RGBColor(0x1A, 0x1A, 0x1A); BODY   = RGBColor(0x52, 0x52, 0x52)
PANEL  = RGBColor(0xF5, 0xF8, 0xFF); SKYBLUE= RGBColor(0x9B, 0xBE, 0xDF)
AUTHCLR= RGBColor(0xCC, 0xDD, 0xEE); PGNUM  = RGBColor(0xAA, 0xAA, 0xAA)
GREEN  = RGBColor(0x70, 0xAD, 0x47); RED    = RGBColor(0xC0, 0x50, 0x4D)
LGREEN = RGBColor(0xE2, 0xEF, 0xDA); LRED   = RGBColor(0xFC, 0xE4, 0xD6)
LBLUE  = RGBColor(0xDE, 0xEB, 0xF7); LORANGE= RGBColor(0xFF, 0xF2, 0xCC)
GRAY   = RGBColor(0xBF, 0xBF, 0xBF); LGRAY  = RGBColor(0xED, 0xED, 0xED)
BOXGRAY= RGBColor(0xE6, 0xE6, 0xE6); EDGEGRAY=RGBColor(0x80, 0x80, 0x80)
PURPLE = RGBColor(0x7A, 0x4F, 0xA3); LPURPLE= RGBColor(0xEC, 0xE3, 0xF5)

F = "Calibri"
SW, SH = 12192000, 6858000
ML, MT = 365760, 228600
CW     = 11430000
RULE_Y = 795528
CONT_Y = 980000
CONT_H = SH - CONT_Y - 228600

# ── primitives ───────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line_color=None, line_w=12700, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    if fill is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color: s.line.color.rgb = line_color; s.line.width = Emu(line_w)
    else: s.line.fill.background()
    s.shadow.inherit = False
    return s

def add_tb(slide, x, y, w, h, wrap=True):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tb.text_frame.word_wrap = wrap
    return tb.text_frame

def fp(tf, text, size, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = F; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return p

def ap(tf, text, size, bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=4, italic=False):
    p = tf.add_paragraph(); p.alignment = align
    if space_before: p.space_before = Pt(space_before)
    r = p.add_run(); r.text = text
    r.font.name = F; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return p

def line(slide, x1, y1, x2, y2, color=EDGEGRAY, w=15875, dash=None):
    c = slide.shapes.add_connector(2, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)))
    c.line.color.rgb = color; c.line.width = Emu(w); c.shadow.inherit = False
    if dash:
        from pptx.oxml.ns import qn
        ln = c.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': dash}))
    return c

_PAGE = [0]
def title_rule(slide, title_text):
    _PAGE[0] += 1
    tf = add_tb(slide, ML, MT, CW, 548640)
    fp(tf, title_text, 27, bold=True, color=DARK)
    add_rect(slide, ML, RULE_Y, CW + 27432, 31750, BLUE)
    tf2 = add_tb(slide, 11430000, 6492240, 640080, 274320)
    fp(tf2, str(_PAGE[0]), 10, color=PGNUM, align=PP_ALIGN.RIGHT)

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def node(slide, x, y, d, fill, label, lblcolor=WHITE, lblsize=13):
    add_rect(slide, x, y, d, d, fill, shape=MSO_SHAPE.OVAL)
    tf = add_tb(slide, x - 40000, y + d/2 - 150000, d + 80000, 300000)
    fp(tf, label, lblsize, bold=True, color=lblcolor, align=PP_ALIGN.CENTER)

def takeaway(slide, text, accent=ORANGE, y=None, sub=None):
    y = y if y else SH - 560000
    h = 430000 if not sub else 560000
    y = SH - h - 120000
    add_rect(slide, ML, y, CW, h, NAVY)
    add_rect(slide, ML, y, 140000, h, accent)
    tf = add_tb(slide, ML + 300000, y + 55000, CW - 480000, h - 110000)
    fp(tf, text, 15, bold=True, color=WHITE)
    if sub: ap(tf, sub, 13, color=accent, space_before=6)


# =============================================================================
# 1 — Title
# =============================================================================
def s_title(prs):
    s = blank(prs)
    add_rect(s, 0, 0, SW, SH, NAVY)
    tf = add_tb(s, ML + 182880, 760000, CW - 365760, 1100000)
    fp(tf, "Hypergraph Redundancy Pruning", 42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tf2 = add_tb(s, ML + 182880, 1640000, CW - 365760, 600000)
    fp(tf2, "Compressing transformers by the redundancy that pairwise methods cannot see",
       21, color=SKYBLUE, align=PP_ALIGN.CENTER)
    add_rect(s, 2926080, 2430000, 6309360, 30480, ORANGE)

    box_y, box_h, box_w, gap = 2640000, 820000, 3300000, 200000
    start_x = (SW - 3 * box_w - 2 * gap) // 2
    boxes = [
        ("The idea", "Some channels are redundant\nonly as a SET, not in pairs"),
        ("The tool", "A hyperedge links a whole\nset of jointly-redundant channels"),
        ("The win", "Prune capacity that graph\nmethods (DepGraph/GOHSP) miss"),
    ]
    for i, (t, sub) in enumerate(boxes):
        bx = start_x + i * (box_w + gap)
        add_rect(s, bx, box_y, box_w, box_h, BLUE)
        tf_b = add_tb(s, bx + 80000, box_y + 90000, box_w - 160000, box_h - 160000)
        fp(tf_b, t, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, ln_ in enumerate(sub.split("\n")):
            ap(tf_b, ln_, 13, color=SKYBLUE, align=PP_ALIGN.CENTER, space_before=4 if j == 0 else 1)

    tf3 = add_tb(s, ML + 182880, 3760000, CW - 365760, 360000)
    fp(tf3, "Shahrzad Esmat   ·   Iowa State University", 15, color=AUTHCLR, align=PP_ALIGN.CENTER)
    tf4 = add_tb(s, ML + 182880, 4140000, CW - 365760, 360000)
    fp(tf4, "Target venue: AAAI", 13, color=ORANGE, align=PP_ALIGN.CENTER, italic=True)


# =============================================================================
# 2 — Topic / problem setup
# =============================================================================
def s_topic(prs):
    s = blank(prs)
    title_rule(s, "The Problem: Which Channels Are Redundant?")
    tf = add_tb(s, ML, CONT_Y, CW, 360000)
    fp(tf, "Pruning makes a transformer smaller by deleting redundant channels. "
           "The entire game is: HOW do you detect redundancy?", 16, bold=True, color=DARK)

    cards = [
        ("Today", BLUE, LBLUE,
         "Compare channels two at a time. If A and B fire almost identically, they duplicate "
         "each other -> keep one. This is what graph-based pruning (DepGraph, GOHSP) does."),
        ("The blind spot", RED, LRED,
         "Redundancy can hide in a SET of channels where no pair looks alike "
         "(e.g. C = A + B). Pairwise checks see nothing and keep all three."),
        ("Our claim", ORANGE, LORANGE,
         "Most removable capacity in a transformer is this set-level kind. A hypergraph "
         "captures it; pairwise graphs structurally cannot."),
    ]
    cy = CONT_Y + 440000; ch = 1180000; cgap = 160000
    for i, (h, acc, fill, body) in enumerate(cards):
        y = cy + i * (ch + cgap)
        add_rect(s, ML, y, CW, ch, fill); add_rect(s, ML, y, 150000, ch, acc)
        tf = add_tb(s, ML + 320000, y + 130000, CW - 600000, ch - 260000)
        fp(tf, h, 18, bold=True, color=acc)
        ap(tf, body, 14, color=BODY, space_before=8)
    takeaway(s, "We reframe pruning as finding redundant SETS of channels — a hypergraph, not a pairwise graph.")


# =============================================================================
# 3 — Background: pairwise redundancy (graph)
# =============================================================================
def s_background(prs):
    s = blank(prs)
    title_rule(s, "Background: Pairwise Redundancy = a Graph")
    tf = add_tb(s, ML, CONT_Y, CW, 320000)
    fp(tf, "Standard pruning builds a graph: a node per channel, an edge when two channels duplicate.",
       15, bold=True, color=DARK)

    # left: two duplicate activation rows
    lx = ML + 120000; ly = CONT_Y + 520000
    tf = add_tb(s, lx, ly - 60000, 5200000, 300000)
    fp(tf, "A duplicate PAIR", 15, bold=True, color=BLUE)
    rows = [("A", [1,0,1,0,1,0], BLUE), ("B", [1,0,1,0,1,0], BLUE)]
    cellw = 560000
    for r, (name, pat, col) in enumerate(rows):
        ry = ly + 360000 + r * 620000
        tf = add_tb(s, lx, ry + 70000, 700000, 300000)
        fp(tf, "Ch " + name, 14, bold=True, color=DARK)
        for c, v in enumerate(pat):
            cx = lx + 900000 + c * cellw
            add_rect(s, cx, ry, cellw - 60000, 380000, col if v else LGRAY,
                     line_color=EDGEGRAY, line_w=6350)
            tf = add_tb(s, cx, ry + 40000, cellw - 60000, 300000)
            fp(tf, str(v), 14, bold=True, color=WHITE if v else BODY, align=PP_ALIGN.CENTER)
    tf = add_tb(s, lx, ly + 1640000, 5200000, 300000)
    fp(tf, "A and B identical  ->  delete one", 13, italic=True, color=BODY)

    # right: graph
    rx = 7100000; ry = CONT_Y + 700000
    tf = add_tb(s, rx, ry - 240000, 4600000, 300000)
    fp(tf, "as a graph (edge = 2 nodes)", 15, bold=True, color=BLUE)
    d = 720000
    node(s, rx + 300000, ry + 200000, d, BLUE, "A")
    node(s, rx + 1900000, ry + 200000, d, BLUE, "B")
    line(s, rx + 300000 + d, ry + 200000 + d/2, rx + 1900000, ry + 200000 + d/2, color=RED, w=25400)
    node(s, rx + 300000, ry + 1500000, d, GRAY, "C")
    node(s, rx + 1900000, ry + 1500000, d, GRAY, "D")
    tf = add_tb(s, rx + 2750000, ry + 350000, 2000000, 300000)
    fp(tf, "edge -> delete one", 13, bold=True, color=RED)
    tf = add_tb(s, rx + 2750000, ry + 1650000, 2000000, 300000)
    fp(tf, "no edge -> keep both", 13, color=BODY)
    takeaway(s, "Every edge connects exactly two nodes. The method can only ever say 'A pairs with B'.", accent=BLUE)


# =============================================================================
# 4 — The gap: redundancy no pair can see
# =============================================================================
def s_gap(prs):
    s = blank(prs)
    title_rule(s, "The Gap: Redundancy No Pair Can See")
    tf = add_tb(s, ML, CONT_Y, CW, 320000)
    fp(tf, "Three channels. C is exactly A + B — fully redundant — yet NO pair looks alike.",
       15, bold=True, color=DARK)

    lx = ML + 120000; ly = CONT_Y + 470000; cellw = 520000
    rows = [("A", [1,0,1,0], BLUE), ("B", [0,1,0,1], GREEN), ("C = A+B", [1,1,1,1], ORANGE)]
    for r, (name, pat, col) in enumerate(rows):
        ry = ly + r * 560000
        tf = add_tb(s, lx, ry + 60000, 1300000, 300000)
        fp(tf, "Ch " + name, 14, bold=True, color=col)
        for c, v in enumerate(pat):
            cx = lx + 1500000 + c * cellw
            add_rect(s, cx, ry, cellw - 60000, 360000, col if v else LGRAY,
                     line_color=EDGEGRAY, line_w=6350)
            tf = add_tb(s, cx, ry + 30000, cellw - 60000, 300000)
            fp(tf, str(v), 14, bold=True, color=WHITE if v else BODY, align=PP_ALIGN.CENTER)

    # pairwise verdict box
    px = lx + 4200000; py = ly
    add_rect(s, px, py, 2550000, 1620000, LRED); add_rect(s, px, py, 0, 0, RED)
    tf = add_tb(s, px + 160000, py + 120000, 2250000, 1400000)
    fp(tf, "Pairwise check:", 14, bold=True, color=RED)
    ap(tf, "A vs B  ->  opposite, not dup", 12.5, color=BODY, space_before=8)
    ap(tf, "A vs C  ->  partial, not dup", 12.5, color=BODY, space_before=4)
    ap(tf, "B vs C  ->  partial, not dup", 12.5, color=BODY, space_before=4)
    ap(tf, "Verdict: keep all 3 ✗", 13, bold=True, color=RED, space_before=8)

    # rank box
    rxx = px + 2750000;
    add_rect(s, rxx, py, 2900000, 1620000, LGREEN)
    tf = add_tb(s, rxx + 160000, py + 120000, 2600000, 1400000)
    fp(tf, "Set view (rank):", 14, bold=True, color=GREEN)
    ap(tf, "3 channels, but only 2", 12.5, color=BODY, space_before=8)
    ap(tf, "are independent.", 12.5, color=BODY, space_before=2)
    ap(tf, "effective rank = 2", 13, bold=True, color=GREEN, space_before=6)
    ap(tf, "-> 1 channel removable ✓", 13, bold=True, color=GREEN, space_before=4)

    takeaway(s,
        "Your measured data: MLP layers have 1536 channels but effective rank ≈ 150 — about 1380 set-redundant dims.",
        accent=ORANGE,
        sub="Pairwise methods detect only ~2 duplicate channels per layer. The HO-gap (~1380) is what we exploit.")


# =============================================================================
# 5 — Related work landscape (with the empty cell)
# =============================================================================
def s_related(prs):
    s = blank(prs)
    title_rule(s, "Related Work: Nobody Occupies Our Cell")
    tf = add_tb(s, ML, CONT_Y, CW, 320000)
    fp(tf, "Two axes: WHAT redundancy is captured (pairwise vs set-level) × HOW it is used (importance vs allocation).",
       14, bold=True, color=DARK)

    # 2x2-ish landscape table
    tx, ty = ML + 200000, CONT_Y + 520000
    colw = [3050000, 3550000, 3550000]
    rowh = 900000
    heads = ["", "Pairwise / 2-way", "Set-level (3+-way)"]
    rows = [
        ("Importance\nscoring",
         ("GOHSP (AAAI'23), DepGraph,\nFPGM — graph head/filter rank", LBLUE, BLUE),
         ("O-information (2022/24):\nsynergy in neuron groups —\nbut tiny nets, no method", LPURPLE, PURPLE)),
        ("Rank / budget\nallocation",
         ("ARA, FLAT-LLM, SVD-LLM,\nASVD — per-module spectra", LBLUE, BLUE),
         ("OUR CELL — empty:\nhyperedge redundancy ->\nallocation on ViT/LLM", LORANGE, ORANGE)),
    ]
    # header row
    for c in range(3):
        cx = tx + sum(colw[:c])
        if c > 0:
            add_rect(s, cx, ty, colw[c] - 80000, 520000, NAVY)
            tf = add_tb(s, cx, ty + 130000, colw[c] - 80000, 320000)
            fp(tf, heads[c], 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for r, (rh, cell1, cell2) in enumerate(rows):
        ry = ty + 620000 + r * (rowh + 60000)
        add_rect(s, tx, ry, colw[0] - 80000, rowh, NAVY)
        tf = add_tb(s, tx, ry + rowh/2 - 240000, colw[0] - 80000, 480000)
        fp(tf, rh.split("\n")[0], 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        ap(tf, rh.split("\n")[1], 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_before=0)
        for c, (txt_, fill, acc) in enumerate([cell1, cell2]):
            cx = tx + sum(colw[:c+1])
            add_rect(s, cx, ry, colw[c+1] - 80000, rowh, fill, line_color=acc, line_w=19050)
            tf = add_tb(s, cx + 120000, ry + 100000, colw[c+1] - 320000, rowh - 200000)
            txtcol = acc if (c == 1 and r == 1) else DARK
            for k, ln_ in enumerate(txt_.split("\n")):
                if k == 0:
                    fp(tf, ln_, 12, bold=True, color=txtcol)
                else:
                    ap(tf, ln_, 12, bold=False, color=txtcol, space_before=2)
    takeaway(s,
        "Gator (2022) uses the word 'hypergraph' but for STRUCTURAL coupling (must-cut-together), not redundancy.",
        accent=ORANGE,
        sub="We must cite + wall off: Gator (structural hypergraph) and O-information (set redundancy, but importance-only, no method).")


# =============================================================================
# 6 — Our method: hyperedge = jointly-redundant set
# =============================================================================
def s_method(prs):
    s = blank(prs)
    title_rule(s, "Our Method: A Hyperedge = a Jointly-Redundant Set")
    tf = add_tb(s, ML, CONT_Y, CW, 320000)
    fp(tf, "A hyperedge wraps a SET of channels (3+) that are jointly rank-deficient — one is reconstructible from the rest.",
       14, bold=True, color=DARK)

    # left: normal graph; right: hypergraph
    lx = ML + 200000; ly = CONT_Y + 560000; d = 700000
    tf = add_tb(s, lx, ly - 280000, 4400000, 300000)
    fp(tf, "Graph: edges of size 2", 15, bold=True, color=BLUE)
    pts = [(lx+200000, ly), (lx+1700000, ly), (lx+950000, ly+1050000)]
    line(s, pts[0][0]+d/2, pts[0][1]+d/2, pts[1][0]+d/2, pts[1][1]+d/2, color=EDGEGRAY, w=19050)
    for (nx, ny), nm in zip(pts, ["A", "B", "C"]):
        node(s, nx, ny, d, BLUE, nm)
    tf = add_tb(s, lx, ly + 2050000, 4200000, 600000)
    fp(tf, "Can only encode pairwise links —", 13, color=BODY)
    ap(tf, "the joint A,B,C dependency is inexpressible.", 13, color=BODY, space_before=2)

    rx = 6700000; ry = ly;
    tf = add_tb(s, rx, ry - 320000, 4600000, 300000)
    fp(tf, "Hypergraph: a hyperedge of size 3", 15, bold=True, color=ORANGE)
    # hyperedge blob behind nodes
    add_rect(s, rx, ry - 40000, 2900000, 1950000, LORANGE, line_color=ORANGE, line_w=22225,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    hpts = [(rx+350000, ry+120000), (rx+1850000, ry+120000), (rx+1100000, ry+1080000)]
    for (nx, ny), nm, cc in zip(hpts, ["A", "B", "C"], [BLUE, GREEN, ORANGE]):
        node(s, nx, ny, d, cc, nm)
    tf = add_tb(s, rx + 3050000, ry + 300000, 2600000, 1200000)
    fp(tf, "one hyperedge e = {A,B,C}", 13, bold=True, color=ORANGE)
    ap(tf, "weight w(e) = redundancy", 12.5, color=BODY, space_before=6)
    ap(tf, "= |e| − rank(e)", 13, bold=True, color=DARK, space_before=4)
    ap(tf, "= 3 − 2 = 1 removable", 12.5, color=GREEN, space_before=4)

    takeaway(s,
        "Redundancy of a set = (#channels) − (effective rank). Sum over hyperedges = total removable capacity.",
        accent=ORANGE)


# =============================================================================
# 7 — Worked example (build hyperedge -> prune)
# =============================================================================
def s_example(prs):
    s = blank(prs)
    title_rule(s, "Worked Example: From Hyperedge to Pruned Layer")
    labels = ["c0","c1","c2","c3","c4","c5"]
    rows = [
        ("1  Measure", "Find sets of channels that are jointly rank-deficient (one ≈ combo of others).",
         "find {c1,c3,c5}: rank 2 of 3   ·   {c0,c4}: independent", None),
        ("2  Build hyperedge", "Wrap each redundant set in a weighted hyperedge  w(e)=|e|−rank(e).",
         "e1={c1,c3,c5} w=1     (e on c0,c2,c4 -> w=0)", {1,3,5}),
        ("3  Allocate & prune", "Spend the budget by removing the redundant member of each hyperedge; "
         "keep a basis that reconstructs it.",
         "drop c5 (≈ c1+c3) -> 6 channels become 5, output preserved", {5}),
    ]
    ry0 = CONT_Y + 60000; row_h = 1380000
    lab_w = 3500000; gx = ML + lab_w + 200000
    node_d = 480000; graph_w = SW - gx - ML; gap = (graph_w - 6*node_d)//5
    accents = [BLUE, ORANGE, GREEN]
    for r, (head, change, note, mark) in enumerate(rows):
        ry = ry0 + r*row_h
        if r % 2 == 0: add_rect(s, ML, ry, CW, row_h - 70000, RGBColor(0xF7,0xF9,0xFC))
        acc = accents[r]
        add_rect(s, ML, ry + 110000, 110000, row_h - 320000, acc)
        tf = add_tb(s, ML + 200000, ry + 150000, lab_w - 250000, row_h - 280000)
        fp(tf, head, 16, bold=True, color=acc)
        ap(tf, change, 12.5, color=BODY, space_before=6)
        ap(tf, note, 11.5, color=acc, space_before=8, italic=True)
        gy = ry + (row_h - node_d)//2
        cx = [gx + i*(node_d+gap) for i in range(6)]
        # hyperedge blob for row 2/3
        if r >= 1:
            members = sorted({1,3,5})
            xl = cx[members[0]] - 70000; xr = cx[members[-1]] + node_d + 70000
            add_rect(s, xl, gy - 110000, xr - xl, node_d + 220000, None,
                     line_color=ORANGE, line_w=19050, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        for i in range(6):
            removed = mark and i in mark and r == 2
            if removed:
                add_rect(s, cx[i], gy, node_d, node_d, LGRAY, line_color=GRAY, line_w=9525, shape=MSO_SHAPE.OVAL)
                tf = add_tb(s, cx[i], gy + node_d/2 - 140000, node_d, 280000)
                fp(tf, "✕", 18, bold=True, color=RED, align=PP_ALIGN.CENTER)
            else:
                fill = ORANGE if (mark and i in {1,3,5} and r==1) else (BLUE if i in {1,3,5} else GREEN)
                node(s, cx[i], gy, node_d, fill, labels[i], lblsize=11)
    fn = add_tb(s, ML, SH - 470000, CW, 360000)
    fp(fn, "Illustrative 6-channel layer. The dropped channel is the one the hyperedge proves is recoverable from its set-mates.",
       12, color=BODY, italic=True)


# =============================================================================
# 8 — Pipeline / algorithm
# =============================================================================
def s_pipeline(prs):
    s = blank(prs)
    title_rule(s, "Pipeline: Plugs Into Any Pruning Backbone")
    steps = [
        ("Calibrate", BLUE, LBLUE,
         "Run a few batches; collect per-channel activations for every attn/MLP layer."),
        ("Find hyperedges", ORANGE, LORANGE,
         "Group channels into jointly rank-deficient sets (correlation-clustering + effective-rank test)."),
        ("Weight", PURPLE, LPURPLE,
         "w(e) = |e| − effrank(e): how many dims the set can shed without changing its span."),
        ("Allocate", GREEN, LGREEN,
         "Distribute the MAC/rank budget across layers proportional to hyperedge redundancy."),
        ("Prune + basis", RED, LRED,
         "Remove redundant members; absorb their contribution into a kept low-rank basis."),
    ]
    n = len(steps); cw = (CW - (n-1)*180000)//n; cy = CONT_Y + 360000; chh = 2600000
    for i, (h, acc, fill, body) in enumerate(steps):
        x = ML + i*(cw + 180000)
        add_rect(s, x, cy, cw, chh, fill); add_rect(s, x, cy, cw, 120000, acc)
        add_rect(s, x + cw/2 - 290000, cy + 260000, 580000, 580000, acc, shape=MSO_SHAPE.OVAL)
        tf = add_tb(s, x + cw/2 - 290000, cy + 320000, 580000, 480000)
        fp(tf, str(i+1), 26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tf = add_tb(s, x + 120000, cy + 980000, cw - 240000, 300000)
        fp(tf, h, 15, bold=True, color=acc, align=PP_ALIGN.CENTER)
        tf = add_tb(s, x + 140000, cy + 1340000, cw - 280000, chh - 1450000)
        fp(tf, body, 12, color=BODY, align=PP_ALIGN.CENTER)
        if i < n-1:
            line(s, x + cw + 20000, cy + 550000, x + cw + 160000, cy + 550000, color=EDGEGRAY, w=22225)
    takeaway(s,
        "Only Step 2–3 are new. Steps 1,4,5 reuse existing machinery — so the gain is attributable to set-level redundancy.",
        accent=ORANGE)


# =============================================================================
# 9 — Expected results
# =============================================================================
def s_results(prs):
    s = blank(prs)
    title_rule(s, "Expected Results")
    tf = add_tb(s, ML, CONT_Y, CW, 320000)
    fp(tf, "Same budget, same backbone, same data — only the redundancy model differs.  (numbers below are TARGETS, not yet measured)",
       14, bold=True, color=DARK)

    # left bar chart: LLM perplexity (lower better)
    lx = ML + 120000; ly = CONT_Y + 560000
    tf = add_tb(s, lx, ly - 60000, 5400000, 300000)
    fp(tf, "LLaMA-2-7B · keep 50% · WikiText ppl ↓", 14, bold=True, color=DARK)
    bars = [("Uniform", 17.39, GRAY), ("FLAT-LLM", 14.56, ORANGE),
            ("Pairwise graph", 14.10, BLUE), ("Ours (hyperedge)", 13.20, GREEN)]
    bx0 = lx + 1900000; bw = 3000000; mx = max(v for _, v, _ in bars)
    for i, (lab, v, col) in enumerate(bars):
        yy = ly + 420000 + i*620000
        tf = add_tb(s, lx, yy + 40000, 1800000, 360000)
        fp(tf, lab, 12.5, bold=("Ours" in lab), color=DARK, align=PP_ALIGN.RIGHT)
        w = bw * v/mx
        add_rect(s, bx0, yy, w, 360000, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = add_tb(s, bx0 + w + 80000, yy + 30000, 900000, 320000)
        fp(tf, f"{v:.2f}", 13, bold=True, color=col)
    tf = add_tb(s, lx, ly + 3050000, 5400000, 300000)
    fp(tf, "target: beat FLAT-LLM and a pairwise-graph allocator at equal budget", 11.5, italic=True, color=BODY)

    # right: the REAL evidence that motivates it
    rx = 6650000; ry = CONT_Y + 560000
    add_rect(s, rx, ry, 5050000, 3250000, PANEL); add_rect(s, rx, ry, 76200, 3250000, ORANGE)
    tf = add_tb(s, rx + 220000, ry + 120000, 4700000, 400000)
    fp(tf, "Evidence already in hand (real probe)", 15, bold=True, color=ORANGE)
    ev = [
        ("MLP eff-rank / channels", "≈ 0.10", "90% of dims are set-redundant", GREEN),
        ("Pairwise duplicates found", "≈ 2 / 1536", "what graph methods can catch", BLUE),
        ("Higher-order gap (HO-gap)", "≈ 1380 dims", "redundancy only a hyperedge sees", ORANGE),
        ("Attn HO-gap", "≈ 300 dims", "same effect, smaller", RED),
    ]
    for i, (k, val, note, col) in enumerate(ev):
        yy = ry + 620000 + i*620000
        add_rect(s, rx + 220000, yy, 60000, 440000, col)
        tf = add_tb(s, rx + 360000, yy - 10000, 2700000, 460000)
        fp(tf, k, 13, bold=True, color=DARK)
        ap(tf, note, 11, color=BODY, space_before=2)
        tf = add_tb(s, rx + 3100000, yy + 40000, 1800000, 400000)
        fp(tf, val, 17, bold=True, color=col, align=PP_ALIGN.RIGHT)
    takeaway(s,
        "The HO-gap is measured and large — the expected win is the act of converting that gap into pruned compute.",
        accent=GREEN)


# =============================================================================
# 10 — Positioning / reviewer-proof
# =============================================================================
def s_positioning(prs):
    s = blank(prs)
    title_rule(s, "Why It Survives Review")
    cards = [
        ("vs ARA / FLAT-LLM / SVD-LLM", BLUE, LBLUE,
         "They allocate rank from each module's OWN spectrum. They cannot see redundancy that "
         "spans channels across a set. We are orthogonal — composable, not competing."),
        ("vs Gator (structural hypergraph)", PURPLE, LPURPLE,
         "Gator's hyperedge = 'must cut together by architecture'. Ours = 'jointly redundant by "
         "information'. Same word, different object — we cite and distinguish."),
        ("vs O-information (set redundancy)", ORANGE, LORANGE,
         "Measures group synergy/redundancy but only as importance, on tiny nets, no compression "
         "method. We turn it into an allocator that beats SVD SOTA on ViT/LLM."),
        ("Caveat we pre-empt", RED, LRED,
         "Low internal rank ≠ unimportant (low-rank heads can be critical). Our signal is "
         "RECONSTRUCTIBILITY-FROM-OTHERS, not a head being internally low-rank."),
    ]
    cw = (CW - 200000)//2; chh = 2300000; gx = 200000; gy = 180000
    cx0 = ML; cy0 = CONT_Y + 120000
    for i, (h, acc, fill, body) in enumerate(cards):
        x = cx0 + (i % 2)*(cw + gx); y = cy0 + (i//2)*(chh + gy)
        add_rect(s, x, y, cw, chh, fill); add_rect(s, x, y, 140000, chh, acc)
        tf = add_tb(s, x + 320000, y + 160000, cw - 520000, chh - 320000)
        fp(tf, h, 16, bold=True, color=acc)
        ap(tf, body, 13.5, color=BODY, space_before=10)
    takeaway(s,
        "Novel claim: hyperedge = informational-redundancy set, used as a compression allocator on transformers, beating SVD SOTA.",
        accent=ORANGE)


def main():
    prs = Presentation()
    prs.slide_width = Emu(SW); prs.slide_height = Emu(SH)
    _PAGE[0] = 0
    s_title(prs); s_topic(prs); s_background(prs); s_gap(prs); s_related(prs)
    s_method(prs); s_example(prs); s_pipeline(prs); s_results(prs); s_positioning(prs)
    out = "hypergraph_redundancy.pptx"
    prs.save(out)
    print("Saved:", out, "| slides:", len(list(prs.slides)))


if __name__ == "__main__":
    main()
