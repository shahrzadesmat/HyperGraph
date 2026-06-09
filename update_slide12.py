from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

P=Presentation("hypergraph_slides.pptx")
# delete current last slide (old slide 12)
sl=P.slides._sldIdLst; sl.remove(list(sl)[-1])

BLANK=P.slide_layouts[6]
DARK=RGBColor(0x1A,0x1A,0x1A); GRAY=RGBColor(0x52,0x52,0x52); MUT=RGBColor(0xAA,0xAA,0xAA)
WHITE=RGBColor(0xFF,0xFF,0xFF); BLUE=RGBColor(0x00,0x68,0xB5); GREEN=RGBColor(0x2E,0x8B,0x57)
GREY=RGBColor(0x9E,0x9E,0x9E); ORANGE=RGBColor(0xFF,0xA3,0x00)
LB=RGBColor(0xDE,0xEB,0xF7); LG=RGBColor(0xE2,0xF0,0xD9); LY=RGBColor(0xFF,0xF2,0xCC)

def rect(s,l,t,w,h,fill,shape=MSO_SHAPE.RECTANGLE):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background(); sp.shadow.inherit=False
    return sp
def txt(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        for (text,sz,bold,col) in line:
            r=p.add_run(); r.text=text; r.font.size=Pt(sz); r.font.bold=bold
            r.font.name="Calibri"; r.font.color.rgb=col
    return tb

s=P.slides.add_slide(BLANK)
txt(s,0.40,0.25,12.5,0.6,[[("Finding 3:  Amplification-Aware Allocation Beats FLAT-LLM",28,True,DARK)]])
rect(s,0.40,0.87,12.53,0.03,BLUE)
txt(s,12.50,7.10,0.70,0.30,[[("12",10,False,MUT)]],align=PP_ALIGN.RIGHT)
txt(s,0.40,1.02,12.5,0.4,[[("Same low-rank primitive, same total budget (",15,True,DARK),
    ("keep 50%",15,True,GREEN),("), Llama-2-7B WikiText perplexity — only the ",15,True,DARK),
    ("allocation",15,True,GREEN),(" differs.",15,True,DARK)]])
# bars
rows=[("No compression",6.87,GRAY),("Uniform",17.39,GREY),
      ("FLAT-LLM  (SOTA)",14.56,ORANGE),("Ours (amplification-joint)",13.96,GREEN)]
x0,y0,rowh,maxw=4.55,1.82,0.70,7.5; mx=max(v for _,v,_ in rows)
for i,(lab,v,col) in enumerate(rows):
    yy=y0+i*rowh; ours="Ours" in lab; sota="FLAT" in lab
    txt(s,0.55,yy+0.06,3.85,0.5,[[(lab,14,True,DARK if (ours or sota) else GRAY)]],align=PP_ALIGN.RIGHT)
    w=maxw*v/mx
    rect(s,x0,yy,w,0.46,col,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s,x0+w+0.1,yy+0.04,1.4,0.4,[[(f"{v:.2f}",15,True,col)]])
txt(s,4.55,y0+4*rowh-0.02,7.5,0.3,[[("perplexity  (lower = better)",11,True,GRAY)]])
# callouts
rect(s,0.40,5.05,6.15,0.92,LG); rect(s,0.40,5.05,0.16,0.92,GREEN)
txt(s,0.75,5.14,5.7,0.8,[[("−4.2%",26,True,GREEN)],[("vs FLAT-LLM  (faithful angular-BI baseline)",12,False,GRAY)]])
rect(s,6.75,5.05,6.15,0.92,LB); rect(s,6.75,5.05,0.16,0.92,BLUE)
txt(s,7.10,5.14,5.7,0.8,[[("−19.7%",26,True,BLUE)],[("vs Uniform allocation",12,False,GRAY)]])
# mechanism / divergence
rect(s,0.40,6.13,12.5,0.95,LY); rect(s,0.40,6.13,0.16,0.95,ORANGE)
txt(s,0.75,6.20,12.0,0.85,[[("Why we win — the signals disagree:  ",12,True,ORANGE),
    ("FLAT-LLM's angular BI protects boundary layers (0, 31, high residual rotation);",12,False,DARK)],
    [("amplification instead protects high-downstream-impact mid-early layers (2–5) and drains the damping tail.  ",12,False,DARK),
    ("Amp captures sensitivity BI misses.",12,True,DARK)]])
P.save("hypergraph_slides.pptx")
print("slide 12 updated; total slides:", len(list(P.slides)))
