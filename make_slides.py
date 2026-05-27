"""
Generate HyperGraph pruning presentation slides
matching the exact style of vdms_slides.pptx.

5 slides:
  1. Problem — ViTs are expensive, need smart pruning
  2. Isomorphic Pruning — existing solution
  3. The Gap — uniform ratio ignores block importance
  4. Our Method — Typed Pruning Hypergraph formula
  5. Example — heterogeneous vs uniform, ablation results
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors (exact from vdms_slides.pptx XML) ─────────────────────────────────
NAVY     = RGBColor(0x1A, 0x27, 0x44)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BLUE     = RGBColor(0x00, 0x68, 0xB5)
ORANGE   = RGBColor(0xFF, 0xA3, 0x00)
DARK     = RGBColor(0x1A, 0x1A, 0x1A)
BODY     = RGBColor(0x52, 0x52, 0x52)
PANEL    = RGBColor(0xF5, 0xF8, 0xFF)
SKYBLUE  = RGBColor(0x9B, 0xBE, 0xDF)
AUTHCLR  = RGBColor(0xCC, 0xDD, 0xEE)
PGNUM    = RGBColor(0xAA, 0xAA, 0xAA)
GREEN    = RGBColor(0x70, 0xAD, 0x47)
RED      = RGBColor(0xC0, 0x50, 0x4D)
LGREEN   = RGBColor(0xE2, 0xEF, 0xDA)
LRED     = RGBColor(0xFC, 0xE4, 0xD6)
LBLUE    = RGBColor(0xDE, 0xEB, 0xF7)
LORANGE  = RGBColor(0xFF, 0xF2, 0xCC)
GRAY     = RGBColor(0xBF, 0xBF, 0xBF)
LGRAY    = RGBColor(0xF2, 0xF2, 0xF2)

F = "Calibri"

# ── Slide geometry (EMU, from vdms_slides.pptx) ───────────────────────────────
SW, SH   = 12192000, 6858000   # slide width / height
ML, MT   = 365760,   228600    # left margin, top margin
CW       = 11430000             # content width
RULE_Y   = 795528               # y of blue underline rule
CONT_Y   = 914400               # y where content starts (below rule)
CONT_H   = SH - CONT_Y - 228600

# ── Primitive helpers ─────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, line_color=None, line_w=12700):
    s = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Emu(line_w)
    else:
        s.line.fill.background()
    return s


def add_tb(slide, x, y, w, h, wrap=True):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tb.text_frame.word_wrap = wrap
    return tb.text_frame


def fp(tf, text, size, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    """Set first paragraph of a fresh text frame."""
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = F
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p


def ap(tf, text, size, bold=False, color=DARK, align=PP_ALIGN.LEFT,
       space_before=4, italic=False):
    """Add a new paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.name = F
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p


def title_rule(slide, title_text, page_num):
    """Standard content slide header: bold title + blue rule + page number."""
    tf = add_tb(slide, ML, MT, CW, 548640)
    fp(tf, title_text, 30, bold=True, color=DARK)
    add_rect(slide, ML, RULE_Y, CW + 27432, 31750, BLUE)
    tf2 = add_tb(slide, 11430000, 6492240, 640080, 274320)
    fp(tf2, str(page_num), 10, color=PGNUM, align=PP_ALIGN.RIGHT)


def left_panel(slide, panel_w=5600000):
    """Add the F5F8FF left panel with blue accent bar. Returns (text_x, text_y)."""
    add_rect(slide, 274320, CONT_Y, panel_w, CONT_H, PANEL)
    add_rect(slide, 274320, CONT_Y, 76200,   CONT_H, BLUE)
    return 457200 + 76200, CONT_Y + 60000   # (text_x, text_y)


def label_box(slide, x, y, w, h, text, size=13, bold=True,
              fill=LBLUE, text_color=BLUE):
    """Small labeled header box (section label style)."""
    add_rect(slide, x, y, w, h, fill)
    tf = add_tb(slide, x + 50000, y + 20000, w - 100000, h - 40000)
    fp(tf, text, size, bold=bold, color=text_color)


# ── Slide factory ─────────────────────────────────────────────────────────────

def blank(prs):
    layout = prs.slide_layouts[6]   # blank layout
    return prs.slides.add_slide(layout)


# =============================================================================
# SLIDE 1 — Title / Problem (dark navy background)
# =============================================================================

def slide1(prs):
    sld = blank(prs)

    # Full dark background
    add_rect(sld, 0, 0, SW, SH, NAVY)

    # Title
    tf = add_tb(sld, ML + 182880, 548640, CW - 365760, 1005840)
    fp(tf, "Efficient Vision Transformers Need Smart Pruning", 40,
       bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    tf2 = add_tb(sld, ML + 182880, 1645560, CW - 365760, 548640)
    fp(tf2, "Not all transformer blocks are equally important",
       22, color=SKYBLUE, align=PP_ALIGN.CENTER)

    # Orange divider
    add_rect(sld, 2926080, 2390000, 6309360, 38100, ORANGE)

    # Three-box problem summary
    box_y = 2530000
    box_h = 685800
    box_w = 3200000
    gap   = 200000
    start_x = (SW - 3 * box_w - 2 * gap) // 2

    boxes = [
        ("DeiT-Small: 4.6G MACs", "Too expensive for edge"),
        ("Structured Pruning", "Remove heads & neurons"),
        ("Key Question", "Which to remove?"),
    ]
    for i, (title, sub) in enumerate(boxes):
        bx = start_x + i * (box_w + gap)
        add_rect(sld, bx, box_y, box_w, box_h, BLUE)
        tf_b = add_tb(sld, bx + 80000, box_y + 80000, box_w - 160000, box_h - 160000)
        fp(tf_b, title, 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        ap(tf_b, sub, 13, color=SKYBLUE, align=PP_ALIGN.CENTER, space_before=6)

    # Author
    tf3 = add_tb(sld, ML + 182880, 3400000, CW - 365760, 365760)
    fp(tf3, "Shahrzad Esmat  |  Iowa State University",
       16, color=AUTHCLR, align=PP_ALIGN.CENTER)

    # Tag line
    tf4 = add_tb(sld, ML + 182880, 3850000, CW - 365760, 365760)
    fp(tf4, "Typed Pruning Hypergraph  ·  ablation on imagenet_10pct",
       14, color=ORANGE, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 2 — How Isomorphic Pruning Works
# =============================================================================

def slide2(prs):
    sld = blank(prs)
    title_rule(sld, "Existing Solution: Isomorphic Pruning", 1)

    # Left panel
    tx, ty = left_panel(sld, panel_w=5600000)

    tf = add_tb(sld, tx, ty, 5000000, 4800000)
    fp(tf, "Uniform Ratio Across All Blocks", 17, bold=True, color=BLUE)
    ap(tf, "", 6)
    ap(tf, "Goal: reduce DeiT-Small from 4.61G → 2.5G MACs", 13, color=BODY, space_before=2)
    ap(tf, "", 5)
    ap(tf, "Formula:", 14, bold=True, color=DARK, space_before=4)
    ap(tf, "   r  =  1 − √( target_MACs / baseline_MACs )", 14,
       bold=True, color=BLUE, italic=True, space_before=2)
    ap(tf, "   r  =  1 − √( 2.5 / 4.61 )  ≈  0.264", 13,
       color=BODY, italic=True, space_before=2)
    ap(tf, "", 5)
    ap(tf, "Applied uniformly to every block:", 14, bold=True, color=DARK, space_before=4)
    ap(tf, "  • Attention:  6 heads → 4 heads (embed 384→256)", 13, color=BODY, space_before=2)
    ap(tf, "  • MLP hidden: 1536 → 1130 neurons", 13, color=BODY, space_before=2)
    ap(tf, "  • All 12 blocks treated identically", 13, color=BODY, space_before=2)
    ap(tf, "", 5)
    ap(tf, "Result:  65.38% top-1  @  2.27G MACs", 14, bold=True, color=GREEN, space_before=4)

    # Right: 12 identical blocks diagram
    rx = 6240000
    ry = CONT_Y + 80000
    bw = 680000
    bh = 310000
    gap = 60000
    label_x = rx + bw + 100000

    # Header
    tf_h = add_tb(sld, rx, ry - 120000, 5500000, 200000)
    fp(tf_h, "All 12 Blocks  — Same Ratio", 14, bold=True, color=DARK)

    for i in range(12):
        by = ry + i * (bh + gap)
        add_rect(sld, rx, by, bw, bh, BLUE)
        tf_b = add_tb(sld, rx + 40000, by + 60000, bw - 80000, bh - 80000)
        fp(tf_b, f"Block {i}", 11, color=WHITE)
        # ratio badge
        add_rect(sld, rx + bw + 80000, by + 60000, 620000, bh - 120000, LBLUE)
        tf_r = add_tb(sld, rx + bw + 80000, by + 60000, 620000, bh - 120000)
        fp(tf_r, "r = 0.264", 11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # Arrow + label
    tf_a = add_tb(sld, rx, ry + 12 * (bh + gap) + 60000, 5000000, 300000)
    fp(tf_a, "↑ Identical pruning ratio for all blocks", 12, color=BODY, italic=True)


# =============================================================================
# SLIDE 3 — The Gap
# =============================================================================

def slide3(prs):
    sld = blank(prs)
    title_rule(sld, "The Gap: Uniform Pruning Ignores Block Importance", 2)

    tx, ty = left_panel(sld, panel_w=5600000)

    tf = add_tb(sld, tx, ty, 5000000, 5000000)
    fp(tf, "Two Critical Observations", 17, bold=True, color=BLUE)
    ap(tf, "", 6)
    ap(tf, "① Block sensitivity spans 3×", 14, bold=True, color=DARK, space_before=4)
    ap(tf, "   Block 0:  S = 1.000  (critical — must keep)", 13, color=BODY, space_before=2)
    ap(tf, "   Block 4:  S = 0.361  (near-redundant)", 13, color=BODY, space_before=2)
    ap(tf, "   → Some blocks could be removed entirely", 13,
       bold=True, color=RED, space_before=2)
    ap(tf, "", 6)
    ap(tf, "② Taylor importance also varies across blocks", 14, bold=True, color=DARK, space_before=4)
    ap(tf, "   High-importance blocks have 2.4× more gradient", 13, color=BODY, space_before=2)
    ap(tf, "   signal than low-importance ones", 13, color=BODY, space_before=2)
    ap(tf, "", 6)
    ap(tf, "③ Isomorphic prunes Block 0 and Block 4 equally", 14, bold=True, color=DARK, space_before=4)
    ap(tf, "   → Wastes MAC budget on redundant blocks", 13, color=RED, space_before=2)
    ap(tf, "   → Over-prunes the most critical block", 13, color=RED, space_before=2)

    # Right: sensitivity bar chart for all 12 blocks
    sensitivities = [1.000, 0.640, 0.470, 0.376, 0.361, 0.372,
                     0.410, 0.427, 0.477, 0.516, 0.567, 0.477]

    rx   = 6240000
    ry   = CONT_Y + 80000
    bw   = 380000
    max_h = 3800000
    gap  = 80000
    base_y = ry + max_h + 200000   # baseline (bars grow upward)

    # Axis label
    tf_yl = add_tb(sld, rx - 200000, ry - 120000, 5500000, 200000)
    fp(tf_yl, "Block Sensitivity  S(i)", 13, bold=True, color=DARK)

    # S_min threshold line
    smin = 0.40
    threshold_y = int(base_y - smin * max_h)
    add_rect(sld, rx - 50000, threshold_y, 12 * (bw + gap) + 100000, 25400, RED)
    tf_thr = add_tb(sld, rx + 12 * (bw + gap) - 200000, threshold_y - 200000, 1500000, 280000)
    fp(tf_thr, "S_min = 0.40", 11, bold=True, color=RED)

    for i, s in enumerate(sensitivities):
        bx = rx + i * (bw + gap)
        bar_h = int(s * max_h)
        by = base_y - bar_h

        # color: green=critical(0), red=redundant(3,4,5), gray=others
        if i == 0:
            fill = GREEN
        elif i in (3, 4, 5):
            fill = RED
        else:
            fill = BLUE

        add_rect(sld, bx, by, bw, bar_h, fill)

        # block label under bar
        tf_l = add_tb(sld, bx, base_y + 40000, bw + gap, 250000)
        fp(tf_l, str(i), 10, color=DARK, align=PP_ALIGN.CENTER)

        # value on top of bar
        tf_v = add_tb(sld, bx, by - 200000, bw + gap, 200000)
        fp(tf_v, f"{s:.2f}", 9, color=DARK, align=PP_ALIGN.CENTER)

    # baseline
    add_rect(sld, rx - 50000, base_y, 12 * (bw + gap) + 100000, 25400, DARK)

    # Legend
    ly = base_y + 350000
    add_rect(sld, rx, ly, 200000, 150000, GREEN)
    tf_lg1 = add_tb(sld, rx + 250000, ly, 1000000, 200000)
    fp(tf_lg1, "Critical", 11, color=BODY)
    add_rect(sld, rx + 1350000, ly, 200000, 150000, RED)
    tf_lg2 = add_tb(sld, rx + 1600000, ly, 1500000, 200000)
    fp(tf_lg2, "Removed (S < 0.40)", 11, color=BODY)
    add_rect(sld, rx + 3200000, ly, 200000, 150000, BLUE)
    tf_lg3 = add_tb(sld, rx + 3450000, ly, 1000000, 200000)
    fp(tf_lg3, "Surviving", 11, color=BODY)


# =============================================================================
# SLIDE 4 — Our Method: Typed Pruning Hypergraph
# =============================================================================

def slide4(prs):
    sld = blank(prs)
    title_rule(sld, "Our Method: Typed Pruning Hypergraph  H = (V′, E_s, E_f)", 3)

    tx, ty = left_panel(sld, panel_w=5600000)

    tf = add_tb(sld, tx, ty, 5000000, 5200000)
    fp(tf, "Three Novel Parameters", 17, bold=True, color=BLUE)

    # S_min
    ap(tf, "", 8)
    ap(tf, "① S_min — Depth Pruning", 15, bold=True, color=DARK, space_before=4)
    ap(tf, "   Remove block i entirely if  S(i) < S_min", 13, color=BLUE,
       italic=True, space_before=2)
    ap(tf, "   Sensitivity: S(i) = mean‖f(x)−f_bypass(x)‖/‖f(x)‖", 12,
       color=BODY, space_before=2)
    ap(tf, "   → Frees MAC budget for surviving blocks", 13, color=BODY, space_before=2)

    # theta
    ap(tf, "", 8)
    ap(tf, "② θ — Heterogeneous Grouping", 15, bold=True, color=DARK, space_before=4)
    ap(tf, "   Group blocks by Taylor importance similarity (BFS)", 13,
       color=BODY, space_before=2)
    ap(tf, "   Per-group ratio:  r_i = r_base × (2 − norm_imp_i)", 13,
       color=BLUE, italic=True, space_before=2)
    ap(tf, "   Important groups pruned less; redundant groups pruned more", 13,
       color=BODY, space_before=2)

    # alpha
    ap(tf, "", 8)
    ap(tf, "③ α — Functional Coupling", 15, bold=True, color=DARK, space_before=4)
    ap(tf, "   I_up(j) = I(j) × (1 + α × Σ w_ij · I(i)/Z)", 13,
       color=BLUE, italic=True, space_before=2)
    ap(tf, "   Predecessors with high importance boost their successors", 13,
       color=BODY, space_before=2)

    # Right: diagram showing the three-step pipeline
    rx = 6240000
    ry = CONT_Y + 60000
    box_w = 5400000
    box_h = 1400000
    gap_h = 200000

    steps = [
        ("① S_min = 0.40", "Remove blocks 3, 4, 5  (S < 0.40)\n3 of 12 blocks removed\nMAC budget redistributed to 9 survivors",
         RED, LRED),
        ("② θ = 0.025", "BFS grouping by Taylor scores → 6 groups\nEach group assigned its own pruning ratio\nImportant blocks (e.g. Block 0) get r ≈ 0.10",
         BLUE, LBLUE),
        ("③ α = 0.3", "Functional edges boost important blocks\nI_up(j) = I(j) × (1 + 0.3 × Σ w_ij·I(i)/Z)\nFactor bounded to [1, 1.3]",
         ORANGE, LORANGE),
    ]

    arrow_x = rx + box_w // 2 - 150000

    for i, (header, body_text, hdr_color, fill_color) in enumerate(steps):
        by = ry + i * (box_h + gap_h)

        # box background
        add_rect(sld, rx, by, box_w, box_h, fill_color)
        # header accent bar
        add_rect(sld, rx, by, 200000, box_h, hdr_color)

        # header text
        tf_h = add_tb(sld, rx + 260000, by + 80000, box_w - 340000, 320000)
        fp(tf_h, header, 16, bold=True, color=hdr_color)

        # body text
        lines = body_text.split('\n')
        tf_b = add_tb(sld, rx + 260000, by + 420000, box_w - 340000, box_h - 480000)
        fp(tf_b, lines[0], 13, color=BODY)
        for line in lines[1:]:
            ap(tf_b, line, 13, color=BODY, space_before=4)

        # arrow between boxes
        if i < 2:
            ay = by + box_h + 30000
            tf_arr = add_tb(sld, arrow_x, ay, 300000, gap_h - 30000)
            fp(tf_arr, "↓", 18, bold=True, color=GRAY, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 5 — Example + Ablation Results
# =============================================================================

def slide5(prs):
    sld = blank(prs)
    title_rule(sld, "Example: Heterogeneous vs Uniform Allocation", 4)

    # ── Left: side-by-side block comparison ───────────────────────────────────
    LX = ML + 50000
    ly_title = CONT_Y + 60000

    tf_hdr = add_tb(sld, LX, ly_title, 5600000, 250000)
    fp(tf_hdr, "Same MAC budget  —  different allocation", 14, bold=True, color=DARK)

    # Sub-headers
    col1_x = LX
    col2_x = LX + 2700000
    col_w  = 2400000
    sub_y  = ly_title + 300000

    tf_c1 = add_tb(sld, col1_x, sub_y, col_w, 250000)
    fp(tf_c1, "Isomorphic", 13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    tf_c2 = add_tb(sld, col2_x, sub_y, col_w, 250000)
    fp(tf_c2, "Ours (H-Pruning)", 13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # Block data: (name, sensitivity, iso_ratio, our_ratio, removed)
    blocks = [
        ("Block 0",  "S=1.00", 0.264, 0.10,  False),
        ("Block 1",  "S=0.64", 0.264, 0.22,  False),
        ("Block 4",  "S=0.36", 0.264, None,   True),   # removed by S_min
        ("Block 9",  "S=0.52", 0.264, 0.30,  False),
    ]

    bh = 680000
    gap = 80000
    max_bw = 1800000
    start_y = sub_y + 350000

    for j, (name, sens, iso_r, our_r, removed) in enumerate(blocks):
        by = start_y + j * (bh + gap)
        mid_y = by + bh // 2 - 130000

        # Name + sensitivity label
        tf_n = add_tb(sld, col1_x, by, 600000, bh)
        fp(tf_n, name, 11, bold=True, color=DARK)
        ap(tf_n, sens, 10, color=BODY, space_before=4)

        # Isomorphic block
        iso_w = int(max_bw * (1.0 - iso_r))
        add_rect(sld, col1_x + 650000, mid_y, iso_w, 260000, BLUE)
        tf_iv = add_tb(sld, col1_x + 650000 + iso_w + 40000, mid_y, 400000, 260000)
        fp(tf_iv, f"r={iso_r:.2f}", 10, color=BLUE)

        # Ours block
        if removed:
            add_rect(sld, col2_x + 50000, mid_y, max_bw, 260000, LRED)
            tf_rem = add_tb(sld, col2_x + 50000 + max_bw // 4, mid_y + 40000,
                            max_bw // 2, 260000)
            fp(tf_rem, "REMOVED", 11, bold=True, color=RED, align=PP_ALIGN.CENTER)
        else:
            our_w = int(max_bw * (1.0 - our_r))
            clr = GREEN if our_r < iso_r else (RED if our_r > iso_r else BLUE)
            add_rect(sld, col2_x + 50000, mid_y, our_w, 260000, clr)
            tf_ov = add_tb(sld, col2_x + 50000 + our_w + 40000, mid_y, 500000, 260000)
            diff = our_r - iso_r
            sign = "+" if diff > 0 else ""
            fp(tf_ov, f"r={our_r:.2f} ({sign}{diff:.2f})", 10, color=clr)

    # note
    note_y = start_y + len(blocks) * (bh + gap) + 80000
    tf_note = add_tb(sld, LX, note_y, 5500000, 250000)
    fp(tf_note, "Budget preserved: Σ r_i × MACs_i  =  Σ r_iso × MACs_i  (by design)",
       11, italic=True, color=BODY)

    # ── Right: ablation results ────────────────────────────────────────────────
    rx = 6400000
    ry = CONT_Y + 60000

    tf_rh = add_tb(sld, rx, ry, 5400000, 300000)
    fp(tf_rh, "Ablation Results  (imagenet_10pct, 20 epochs)", 14, bold=True, color=DARK)
    tf_rs = add_tb(sld, rx, ry + 310000, 5400000, 200000)
    fp(tf_rs, "DeiT-Small, target 2.5G MACs, seed=42", 12, color=BODY, italic=True)

    # Bar chart of ablation
    steps = [
        ("iso_baseline\n(S_min=0, θ=1.0, α=0)", 0.6538, GRAY),
        ("+ S_min=0.40\n(depth pruning)",         0.6562, BLUE),
        ("+ θ=0.025\n(heterogeneous groups)",      0.6594, ORANGE),
        ("+ α=0.3\n(functional coupling)",          0.6624, GREEN),
    ]

    bar_start_y = ry + 620000
    bar_h = 560000
    bar_gap = 100000
    bar_max_w = 3600000
    min_acc = 0.645
    acc_range = 0.025   # 0.645 to 0.670

    for i, (label, acc, color) in enumerate(steps):
        by = bar_start_y + i * (bar_h + bar_gap)

        # bar
        bar_w = int(bar_max_w * (acc - min_acc) / acc_range)
        bar_w = max(bar_w, 200000)
        add_rect(sld, rx, by + 100000, bar_w, bar_h - 200000, color)

        # accuracy label
        tf_acc = add_tb(sld, rx + bar_w + 60000, by + 100000, 900000, bar_h - 200000)
        fp(tf_acc, f"{acc:.2%}", 14, bold=True, color=color)

        # delta label for rows > 0
        if i > 0:
            delta = acc - steps[i-1][1]
            tf_d = add_tb(sld, rx + bar_w + 60000 + 850000, by + 100000, 700000, bar_h - 200000)
            fp(tf_d, f"(+{delta:.2%})", 12, color=GREEN, bold=True)

        # step label
        tf_lbl = add_tb(sld, rx + bar_max_w + 200000, by, 1900000, bar_h)
        lines = label.split('\n')
        fp(tf_lbl, lines[0], 12, bold=True, color=DARK)
        if len(lines) > 1:
            ap(tf_lbl, lines[1], 11, color=BODY, space_before=3)

    # Total gain callout
    cy = bar_start_y + 4 * (bar_h + bar_gap) + 80000
    add_rect(sld, rx, cy, 5400000, 380000, LGREEN)
    add_rect(sld, rx, cy, 120000, 380000, GREEN)
    tf_tot = add_tb(sld, rx + 180000, cy + 60000, 5100000, 300000)
    fp(tf_tot, "Total gain over isomorphic baseline:  +0.86pp  at same 2.27G MACs",
       13, bold=True, color=GREEN)


# =============================================================================
# Main
# =============================================================================

def main():
    prs = Presentation()
    prs.slide_width  = Emu(SW)
    prs.slide_height = Emu(SH)

    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)

    out = "/work/hdd/bdjd/hypergraph_pruning/hypergraph_slides.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
