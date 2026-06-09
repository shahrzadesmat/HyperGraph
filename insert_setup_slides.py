from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

P=Presentation("hypergraph_slides.pptx")
BLANK=P.slide_layouts[6]
DARK=RGBColor(0x1A,0x1A,0x1A); GRAY=RGBColor(0x52,0x52,0x52); MUT=RGBColor(0xAA,0xAA,0xAA)
WHITE=RGBColor(0xFF,0xFF,0xFF); BLUE=RGBColor(0x00,0x68,0xB5); RED=RGBColor(0xC0,0x50,0x4D)
ORANGE=RGBColor(0xFF,0xA3,0x00); GREEN=RGBColor(0x2E,0x8B,0x57)
LB=RGBColor(0xDE,0xEB,0xF7); LY=RGBColor(0xFF,0xF2,0xCC); LP=RGBColor(0xFC,0xE4,0xD6); LG=RGBColor(0xE2,0xF0,0xD9)

def rect(s,l,t,w,h,fill,shape=MSO_SHAPE.RECTANGLE):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        for (text,sz,bold,col) in line:
            r=p.add_run(); r.text=text; r.font.size=Pt(sz); r.font.bold=bold; r.font.name="Calibri"; r.font.color.rgb=col
    return tb
def header(s,title):
    txt(s,0.40,0.25,12.5,0.6,[[(title,28,True,DARK)]]); rect(s,0.40,0.87,12.53,0.03,BLUE)
    txt(s,12.50,7.10,0.70,0.30,[[("",10,False,MUT)]],align=PP_ALIGN.RIGHT)  # placeholder, renumbered later

# ---------- NEW SLIDE A : what SOTA / FLAT-LLM does ----------
sA=P.slides.add_slide(BLANK); header(sA,"Background:  How SOTA Compresses an LLM (e.g. FLAT-LLM)")
txt(sA,0.40,1.02,12.5,0.4,[[("To shrink a model, low-rank methods replace each layer's FFN with a ",15,True,DARK),
    ("smaller approximation",15,True,BLUE),(" that keeps only its most useful directions.",15,True,DARK)]])
cards=[("1","The primitive",BLUE,LB,
        "Each FFN is approximated by a low-rank factor in activation space (keep the top-k PCA directions of its hidden units).  Fewer directions = less compute."),
       ("2","The budget",ORANGE,LY,
        "A fixed total rank is shared across all 32 layers.  The core question: how much rank does each layer get?"),
       ("3","The allocation (FLAT-LLM)",RED,LP,
        "Score each layer by a LOCAL importance — FLAT-LLM uses angular Block-Influence: how much that block rotates the residual.  Give rank in proportion (IPRS).")]
y=1.62
for n,h,acc,fill,body in cards:
    rect(sA,0.40,y,12.5,1.36,fill); rect(sA,0.40,y,0.16,1.36,acc)
    rect(sA,0.73,y+0.31,0.74,0.74,acc,shape=MSO_SHAPE.OVAL)
    txt(sA,0.73,y+0.34,0.74,0.7,[[(n,26,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(sA,1.69,y+0.18,10.9,0.4,[[(h,18,True,acc)]])
    txt(sA,1.69,y+0.62,10.9,0.7,[[(body,13,False,GRAY)]]); y+=1.50
rect(sA,0.40,6.45,12.5,0.72,LB); rect(sA,0.40,6.45,0.16,0.72,BLUE)
txt(sA,0.75,6.52,12.0,0.6,[[("All SOTA low-rank methods (ASVD, SVD-LLM, FLAT-LLM) share this template:  ",13,True,BLUE),
    ("a per-layer LOCAL score, and each layer allocated independently.",13,False,DARK)]])

# ---------- NEW SLIDE B : the weakness we cure ----------
sB=P.slides.add_slide(BLANK); header(sB,"The Weakness We Cure:  Layers Are Scored in Isolation")
txt(sB,0.40,1.02,12.5,0.4,[[("A local, per-layer score cannot see how a layer's compression error ",15,True,DARK),
    ("travels through the rest of the network.",15,True,RED)]])
# problem (left, red) / cure (right, green)
rect(sB,0.40,1.62,6.15,3.95,LP); rect(sB,0.40,1.62,0.16,3.95,RED)
txt(sB,0.75,1.78,5.6,0.5,[[("The blind spot",20,True,RED)]])
txt(sB,0.75,2.40,5.6,3.0,[[("Transformer layers form a chain through the residual stream.",13,True,DARK)],
    [("",6,False,DARK)],
    [("An error introduced when compressing an ",13,False,GRAY),("early",13,True,DARK),
     (" layer does not stay local — it ",13,False,GRAY),("propagates and amplifies",13,True,RED),
     (" downstream (we measure up to ~30× by the output).",13,False,GRAY)],
    [("",6,False,DARK)],
    [("A local score (reconstruction error, block-influence) is computed at ",13,False,GRAY),
     ("one layer in isolation",13,True,DARK),(" — it is structurally blind to this propagation.",13,False,GRAY)]])
rect(sB,6.75,1.62,6.15,3.95,LG); rect(sB,6.75,1.62,0.16,3.95,GREEN)
txt(sB,7.10,1.78,5.6,0.5,[[("Our cure  (the novelty)",20,True,GREEN)]])
txt(sB,7.10,2.40,5.6,3.0,[[("Allocate rank by ",13,False,GRAY),("downstream impact",13,True,DARK),
     (", not local importance.",13,False,GRAY)],
    [("",6,False,DARK)],
    [("• Measure each layer's ",13,False,GRAY),("error amplification",13,True,GREEN),
     (" to the final output.",13,False,GRAY)],
    [("• Account for ",13,False,GRAY),("cross-layer coupling",13,True,GREEN),
     (" (errors interact; allocate jointly, not greedily).",13,False,GRAY)],
    [("",6,False,DARK)],
    [("Protect layers whose errors blow up;  compress hard where they damp.",13,True,DARK)]])
rect(sB,0.40,5.75,12.5,1.05,LB); rect(sB,0.40,5.75,0.16,1.05,BLUE)
txt(sB,0.75,5.85,12.0,0.9,[[("Novelty in one line:  ",13,True,BLUE),
    ("the first error-propagation-aware, cross-layer-coupled rank allocation for low-rank LLM compression.",13,True,DARK)],
    [("Everything prior allocates per layer in isolation; we allocate over the propagation graph.  ",12,False,GRAY),
    ("(Result on the next slides: beats FLAT-LLM at the same budget.)",12,False,GREEN)]])

# ---------- reorder: move the 2 new slides (now last) to before old Finding-1 (index 9) ----------
sldIdLst=P.slides._sldIdLst
ids=list(sldIdLst)            # 14 entries; new A=ids[12], B=ids[13]
newA, newB = ids[12], ids[13]
ref = ids[9]                  # old "Finding 1" slide
sldIdLst.remove(newA); sldIdLst.remove(newB)
ref.addprevious(newA); ref.addprevious(newB)   # -> ...,8, A, B, Finding1, ...

# ---------- renumber page numbers on shifted content slides (final indices 9..13) ----------
def numbox(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.left and Emu(sh.left).inches>11 and Emu(sh.top).inches>6.8:
            return sh
    return None
slides=list(P.slides)
for idx in range(9,14):
    nb=numbox(slides[idx])
    if nb is None: continue
    tf=nb.text_frame
    if tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text=str(idx)
    else:
        r=tf.paragraphs[0].add_run(); r.text=str(idx); r.font.size=Pt(10); r.font.name="Calibri"; r.font.color.rgb=MUT
        tf.paragraphs[0].alignment=PP_ALIGN.RIGHT

P.save("hypergraph_slides.pptx")
print("done. total slides:", len(list(P.slides)))
for i,s in enumerate(P.slides):
    t=[sh.text_frame.text for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
    print(f"  {i}: {t[0][:55] if t else '?'}")
