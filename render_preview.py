import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle, FancyBboxPatch
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

P=Presentation("hypergraph_slides.pptx")
SW,SH=13.3333,7.5
def rgb(sh):
    try:
        if sh.fill.type is not None and sh.fill.fore_color and sh.fill.fore_color.type is not None:
            c=sh.fill.fore_color.rgb; return (c[0]/255,c[1]/255,c[2]/255)
    except: pass
    return None
def trgb(run):
    try:
        if run.font.color and run.font.color.type is not None:
            c=run.font.color.rgb; return (c[0]/255,c[1]/255,c[2]/255)
    except: pass
    return (0,0,0)

idxs=[9,13]
fig,axes=plt.subplots(len(idxs),1,figsize=(13.33,15))
axes=[axes] if len(idxs)==1 else axes
for ax,idx in zip(axes,idxs):
    s=P.slides[idx]
    ax.set_xlim(0,SW); ax.set_ylim(0,SH); ax.invert_yaxis(); ax.set_aspect('equal')
    ax.add_patch(Rectangle((0,0),SW,SH,fill=True,facecolor='white',edgecolor='black',lw=2))
    ax.set_xticks([]); ax.set_yticks([]); ax.set_title(f"slide {idx}",fontsize=9,loc='left')
    for sh in s.shapes:
        l,t,w,h=Emu(sh.left).inches,Emu(sh.top).inches,Emu(sh.width).inches,Emu(sh.height).inches
        fc=rgb(sh)
        if fc is not None:
            ax.add_patch(Rectangle((l,t),w,h,facecolor=fc,edgecolor='none',alpha=0.95))
        if sh.has_text_frame and sh.text_frame.text.strip():
            paras=list(sh.text_frame.paragraphs)
            for pi,p in enumerate(paras):
                if not p.runs: continue
                full="".join(r.text for r in p.runs)
                if not full.strip(): continue
                r0=p.runs[0]; sz=(r0.font.size.pt if r0.font.size else 12)
                al={1:'left',2:'center',3:'right'}.get(p.alignment,'left') if p.alignment else 'left'
                px=l if al!='center' else l+w/2
                if al=='right': px=l+w
                yi=t+0.16+pi*(sz/52.0)
                ax.text(px, yi, full[:95], fontsize=min(sz*0.6,14), color=trgb(r0),
                        ha=al, va='top', fontweight=('bold' if r0.font.bold else 'normal'))
plt.tight_layout()
plt.savefig("deck_preview.png",dpi=130,bbox_inches='tight')
print("saved deck_preview.png")
