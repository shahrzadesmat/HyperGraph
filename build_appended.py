import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

P=Presentation("hypergraph_slides.pptx")
BLANK=P.slide_layouts[6]
SW,SH=13.3333,7.5
DARK=RGBColor(0x1A,0x1A,0x1A); GRAY=RGBColor(0x52,0x52,0x52); MUT=RGBColor(0xAA,0xAA,0xAA)
WHITE=RGBColor(0xFF,0xFF,0xFF); BLUE=RGBColor(0x00,0x68,0xB5); RED=RGBColor(0xC0,0x50,0x4D)
ORANGE=RGBColor(0xFF,0xA3,0x00); GREEN=RGBColor(0x2E,0x8B,0x57); GREY=RGBColor(0x9E,0x9E,0x9E)
LB=RGBColor(0xDE,0xEB,0xF7); LY=RGBColor(0xFF,0xF2,0xCC); LP=RGBColor(0xFC,0xE4,0xD6); LG=RGBColor(0xE2,0xF0,0xD9)

BOXES=[]  # (slide_idx, name, l,t,w,h) for validation
def rect(s,si,l,t,w,h,fill,shape=MSO_SHAPE.RECTANGLE,name="rect"):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background(); sp.shadow.inherit=False
    BOXES.append((si,name,l,t,w,h)); return sp
def txt(s,si,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,name="txt"):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        for (text,sz,bold,col) in line:
            r=p.add_run(); r.text=text; r.font.size=Pt(sz); r.font.bold=bold; r.font.name="Calibri"; r.font.color.rgb=col
    BOXES.append((si,name,l,t,w,h)); return tb
def base(si,title):
    s=P.slides.add_slide(BLANK)
    txt(s,si,0.40,0.25,12.5,0.6,[[(title,26,True,DARK)]],name="title")
    rect(s,si,0.40,0.87,12.53,0.03,BLUE,name="rule")
    txt(s,si,12.50,7.10,0.70,0.30,[[(str(si),10,False,MUT)]],align=PP_ALIGN.RIGHT,name="pagenum")
    return s

# ===== SLIDE 9 : background + what is rank =====
si=9; s=base(si,"Background:  How SOTA Compresses an LLM (e.g. FLAT-LLM)")
txt(s,si,0.40,1.00,12.5,0.4,[[("Low-rank methods shrink a model by replacing each layer with a ",14,True,DARK),
    ("smaller approximation",14,True,BLUE),(" that keeps only its most useful patterns.",14,True,DARK)]])
# RANK explainer banner
rect(s,si,0.40,1.50,12.5,1.18,LB,name="rankbanner"); rect(s,si,0.40,1.50,0.16,1.18,BLUE,name="rankbar")
txt(s,si,0.75,1.58,12.0,1.05,[[("What is “rank”, simply?    ",15,True,BLUE),
    ("A layer does its job using a number of independent patterns.  ",13,False,DARK),
    ("“Rank” = how many of those patterns we keep.",13,True,DARK)],
    [("Keeping fewer patterns (lower rank) makes the layer smaller and faster — like describing a photo with a few broad strokes instead of every pixel.   ",13,False,GRAY),
     ("So “how much to compress a layer” = “how many patterns (rank) to keep.”",13,True,DARK)]])
# 3 columns
cols=[("1","The primitive",BLUE,LB,"Approximate each FFN by a low-rank factor — keep its top patterns (directions). Fewer kept = less compute."),
      ("2","The budget",ORANGE,LY,"A fixed TOTAL rank is shared across all 32 layers. The core question: how much rank does each layer get?"),
      ("3","Allocation — FLAT-LLM",RED,LP,"Score each layer by its block-influence — how much it changes the residual (the angle between what goes in and what comes out). More influence → keep more rank.")]
cw,gap,x0,cy,ch=4.03,0.205,0.40,2.86,2.42
for i,(n,h,acc,fill,body) in enumerate(cols):
    cx=x0+i*(cw+gap)
    rect(s,si,cx,cy,cw,ch,fill,name="col"); rect(s,si,cx,cy,cw,0.10,acc,name="coltop")
    rect(s,si,cx+0.22,cy+0.28,0.62,0.62,acc,shape=MSO_SHAPE.OVAL,name="circ")
    txt(s,si,cx+0.22,cy+0.30,0.62,0.58,[[(n,24,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,name="cn")
    txt(s,si,cx+0.98,cy+0.30,cw-1.1,0.6,[[(h,15,True,acc)]],anchor=MSO_ANCHOR.MIDDLE,name="ch")
    txt(s,si,cx+0.22,cy+1.08,cw-0.44,1.2,[[(body,12,False,GRAY)]],name="cb")
rect(s,si,0.40,5.55,12.5,0.80,LB,name="botstrip"); rect(s,si,0.40,5.55,0.16,0.80,BLUE,name="botbar")
txt(s,si,0.75,5.64,12.0,0.65,[[("All SOTA low-rank methods (ASVD, SVD-LLM, FLAT-LLM) share this template:  ",13,True,BLUE),
    ("a per-layer LOCAL score, with each layer allocated independently.",13,False,DARK)]],anchor=MSO_ANCHOR.MIDDLE)

# ===== SLIDE 10 : low-rank vs other structured pruning + shared weakness =====
si=10; s=base(si,"Low-Rank vs. Other Structured Pruning")
txt(s,si,0.40,1.00,12.5,0.4,[[("Both are “structured” pruning (regular, hardware-friendly shapes) — they differ in ",14,True,DARK),("what they keep.",14,True,BLUE)]])
# two contrast cards
rect(s,si,0.40,1.50,6.15,2.02,LP,name="cl"); rect(s,si,0.40,1.50,0.16,2.02,ORANGE,name="clb")
txt(s,si,0.75,1.63,5.6,0.4,[[("Remove units",18,True,ORANGE)]],name="clh")
txt(s,si,0.75,2.10,5.6,1.0,[[("Delete whole ",13,False,GRAY),("channels, attention heads, or blocks",13,True,DARK),(".  The units left behind are the originals, unchanged.",13,False,GRAY)]],name="clb1")
txt(s,si,0.75,3.04,5.6,0.4,[[("‹ the ViT channel / block pruning earlier in this talk ›",12,False,ORANGE)]],name="clb2")
rect(s,si,6.75,1.50,6.15,2.02,LB,name="cr"); rect(s,si,6.75,1.50,0.16,2.02,BLUE,name="crb")
txt(s,si,7.10,1.63,5.6,0.4,[[("Reduce rank   (these methods)",18,True,BLUE)]],name="crh")
txt(s,si,7.10,2.10,5.6,1.3,[[("Keep no clean subset.  Replace a weight matrix  ",13,False,GRAY),("W ≈ A · B",13,True,DARK),("  with two smaller matrices — keeping a low-dimensional ",13,False,GRAY),("rotated subspace",13,True,DARK),(" (mixtures of neurons).   “Rank” = the inner size.",13,False,GRAY)]],name="crb1")
# shared weakness
rect(s,si,0.40,3.68,12.5,1.32,LP,name="wk"); rect(s,si,0.40,3.68,0.16,1.32,RED,name="wkb")
txt(s,si,0.75,3.78,12.0,1.18,[[("The weakness both share:  ",14,True,RED),("either way, the method decides ",13,False,DARK),("per layer how much to cut",13,True,DARK),(", from a ",13,False,DARK),("LOCAL",13,True,DARK),(" score (importance / reconstruction at that layer alone).",13,False,DARK)],
    [("It is blind to how an ",13,False,GRAY),("early layer's error amplifies downstream",13,True,RED),(" (~30× by the output) — the chain is ignored.",13,False,GRAY)]],name="wkt")
# cure
rect(s,si,0.40,5.16,12.5,1.45,LG,name="cu"); rect(s,si,0.40,5.16,0.16,1.45,GREEN,name="cub")
txt(s,si,0.75,5.26,12.0,1.3,[[("Our approach:  ",14,True,GREEN),("allocate rank by ",13,False,DARK),("downstream amplification",13,True,GREEN),(" and ",13,False,DARK),("cross-layer coupling",13,True,GREEN),(", not a purely local per-layer score.",13,False,DARK)],
    [("Protect layers whose errors blow up; compress hard where they damp.   ",13,False,GRAY),("Empirically: beats FLAT-LLM at the same budget (next slide).",13,True,DARK)]],name="cut")

# ===== SLIDE 11 : Finding 1 — FIXED chart =====
si=11; s=base(si,"Finding 1:  Compression Error Amplifies with Depth")
txt(s,si,0.40,1.00,12.5,0.4,[[("A layer's ",14,True,DARK),("local",14,True,RED),
    (" error is not its true cost — it propagates to the output, and the amplification depends on depth.",14,True,DARK)]])
rect(s,si,0.40,1.52,12.5,0.56,LB,name="probe")
txt(s,si,0.62,1.59,12.1,0.45,[[("Probe:  ",12,True,BLUE),("rank-truncate ONE layer's FFN, measure  ",12,False,GRAY),
    ("Δ output ÷ Δ local",12,True,DARK),("  = how much its error blows up by the time it reaches the logits.",12,False,GRAY)]],anchor=MSO_ANCHOR.MIDDLE)
# chart band: bars grow UP from baseline=5.05, max height 2.55 (top>=2.50, clear of probe strip end 2.08)
data=[("L0",30.8),("L1",14.9),("L3",13.4),("L5",5.4),("L10",2.0),("L16",1.8),("L24",1.2),("L31",0.7)]
x0,bw,gap,bottom,maxh=0.95,1.05,0.40,5.05,2.55
mx=math.log10(max(v for _,v in data)+1)
for i,(lab,v) in enumerate(data):
    h=maxh*math.log10(v+1)/mx; cx=x0+i*(bw+gap); top=bottom-h
    col=RED if v>=5 else (ORANGE if v>=1.3 else BLUE)
    rect(s,si,cx,top,bw,h,col,shape=MSO_SHAPE.ROUNDED_RECTANGLE,name=f"bar{lab}")
    txt(s,si,cx-0.20,top-0.32,bw+0.4,0.3,[[(f"{v:.1f}×",13,True,col)]],align=PP_ALIGN.CENTER,name=f"val{lab}")
    txt(s,si,cx-0.20,bottom+0.06,bw+0.4,0.26,[[(lab,12,False,GRAY)]],align=PP_ALIGN.CENTER,name=f"x{lab}")
txt(s,si,0.62,5.42,7.0,0.26,[[("amplification = Δe2e ÷ Δlocal   (log scale; L0 = nearest input)",11,True,GRAY)]],name="axis")
txt(s,si,9.0,5.42,3.9,0.26,[[("early → amplify ~30×    ",11,True,RED),("late → damp <1×",11,True,BLUE)]],align=PP_ALIGN.RIGHT,name="annot")
rect(s,si,0.40,6.35,12.5,0.78,LP,name="impl"); rect(s,si,0.40,6.35,0.16,0.78,RED,name="implb")
txt(s,si,0.75,6.43,12.0,0.64,[[("Implication:  ",13,True,RED),
    ("the true cost of compressing a layer varies ~43× by depth.  Methods that allocate rank by ",13,False,DARK),
    ("local",13,True,DARK),(" error (FLAT-LLM, ASVD, SVD-LLM) optimize the wrong quantity.",13,False,DARK)]],anchor=MSO_ANCHOR.MIDDLE)

# ===== SLIDE 12 : Finding 2 — banded graph =====
si=12; s=base(si,"Finding 2:  Layer Errors Couple — a Graph, Not Just a Curve")
txt(s,si,0.40,1.00,12.5,0.4,[[("Do layers' compression errors ",14,True,DARK),("interact",14,True,BLUE),
    ("?  If so, rank must be allocated jointly, not greedily per-layer.",14,True,DARK)]])
cards=[("ι","Errors are non-additive",ORANGE,LY,
        "Compressing two layers together ≠ sum of each alone.  Interaction (mean 0.60) is large — joint effects matter."),
       ("|cos|","Coupling is local (banded)",BLUE,LB,
        "Error directions align for nearby layers, decay with distance:  adjacent 0.46  ›  distant 0.28.  A banded graph."),
       ("28–31","The tail decouples",RED,LP,
        "The final block is weakly coupled to all others (|cos| 0.12–0.24) — a separate, compressible group.")]
y=1.58
for big,h,acc,fill,body in cards:
    rect(s,si,0.40,y,12.5,1.30,fill,name="card"); rect(s,si,0.40,y,0.16,1.30,acc,name="cardb")
    rect(s,si,0.73,y+0.30,0.72,0.72,acc,shape=MSO_SHAPE.OVAL,name="circ")
    txt(s,si,0.66,y+0.33,0.86,0.66,[[(big,17 if len(big)>2 else 24,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,name="big")
    txt(s,si,1.69,y+0.16,10.9,0.4,[[(h,17,True,acc)]],name="h")
    txt(s,si,1.69,y+0.60,10.9,0.66,[[(body,13,False,GRAY)]],name="b")
    y+=1.44
rect(s,si,0.40,6.35,12.5,0.78,LB,name="bot"); rect(s,si,0.40,6.35,0.16,0.78,BLUE,name="botb")
txt(s,si,0.75,6.43,12.0,0.64,[[("The adjacency is real:  ",13,True,BLUE),
    ("a banded error-coupling graph over the residual stream → rank allocation is a ",13,False,DARK),
    ("joint",13,True,DARK),(" problem, which per-layer greedy methods cannot solve.",13,False,DARK)]],anchor=MSO_ANCHOR.MIDDLE)

# ===== SLIDE 13 (new) : how the three methods allocate rank =====
si=13; s=base(si,"How the Three Methods Allocate the Rank Budget")
txt(s,si,0.40,1.00,12.5,0.4,[[("Every method spends the ",14,True,DARK),("same total rank",14,True,GREEN),
    (" — they differ only in ",14,True,DARK),("which layers get more.",14,True,BLUE)]])
LGREY=RGBColor(0xEE,0xEE,0xEE)
cols=[("Uniform",GREY,LGREY,"Give every layer the same rank.",
       "Uses no signal at all — the naive baseline.",[0.5]*8),
      ("FLAT-LLM  (SOTA)",ORANGE,LY,"Rank set by block-influence (how much a layer rotates the residual).",
       "Protects boundary layers (0 & 31); compresses the middle.",[1.0,0.62,0.62,0.55,0.51,0.29,0.39,0.93]),
      ("Ours — amplification-joint",GREEN,LG,"Rank set by measured downstream amplification, allocated jointly across coupled layers.",
       "Protects high-impact early / mid layers (2–5); drains the damping tail.",[0.57,0.90,0.90,0.52,0.45,0.40,0.37,0.37])]
labs=["0","2","5","10","16","24","30","31"]
cw,gap,x0,cy,ch=4.03,0.205,0.40,1.52,4.85
for i,(name,acc,fill,rule,protect,prof) in enumerate(cols):
    cx=x0+i*(cw+gap)
    rect(s,si,cx,cy,cw,ch,fill,name="col"); rect(s,si,cx,cy,cw,0.10,acc,name="ct")
    txt(s,si,cx+0.22,cy+0.22,cw-0.44,0.4,[[(name,15,True,acc)]],name="cn")
    txt(s,si,cx+0.22,cy+0.72,cw-0.44,1.0,[[(rule,12,False,GRAY)]],name="cr")
    bb,bmaxh=cy+3.40,1.40
    nb=len(prof); bw=0.30; gp=((cw-0.60)-nb*bw)/(nb-1); bx0=cx+0.30
    for j,hv in enumerate(prof):
        bh=bmaxh*hv; bxx=bx0+j*(bw+gp)
        rect(s,si,bxx,bb-bh,bw,bh,acc,shape=MSO_SHAPE.ROUNDED_RECTANGLE,name="pb")
        txt(s,si,bxx-0.10,bb+0.03,bw+0.2,0.2,[[(labs[j],7,False,GRAY)]],align=PP_ALIGN.CENTER,name="pl")
    txt(s,si,cx+0.22,cy+3.74,cw-0.44,0.3,[[("rank per layer  (L0 → L31)",9,False,GRAY)]],name="pcap")
    txt(s,si,cx+0.22,cy+4.06,cw-0.44,0.7,[[("Protects:  ",11,True,acc),(protect,11,False,GRAY)]],name="cp")
txt(s,si,0.40,6.55,12.5,0.35,[[("Same budget, different shape — next slide: which shape gives the lowest perplexity.",12,True,DARK)]],align=PP_ALIGN.CENTER,name="foot")

# ===== SLIDE 14 : Finding 3 — beats FLAT-LLM =====
si=14; s=base(si,"Finding 3:  Amplification-Aware Allocation Beats FLAT-LLM")
txt(s,si,0.40,1.00,12.5,0.4,[[("Same low-rank primitive, same total budget (",14,True,DARK),
    ("keep 50%",14,True,GREEN),("), Llama-2-7B WikiText perplexity — only the ",14,True,DARK),
    ("allocation",14,True,GREEN),(" differs.",14,True,DARK)]])
rows=[("No compression","the full, uncompressed model",6.87,GREY),
      ("Uniform","every layer gets the SAME rank",17.39,GREY),
      ("FLAT-LLM (SOTA)","rank set by block-influence score",14.56,ORANGE),
      ("Ours","more rank to high-amplification layers, jointly",13.96,GREEN)]
x0,y0,rowh,maxw=4.95,1.80,0.72,6.6; mx=max(v for _,_,v,_ in rows)
for i,(lab,gloss,v,col) in enumerate(rows):
    yy=y0+i*rowh; hot=("Ours" in lab or "FLAT" in lab)
    txt(s,si,0.30,yy-0.05,4.45,0.62,[[(lab,14,True,DARK if hot else GRAY)],[(gloss,10,False,GRAY)]],align=PP_ALIGN.RIGHT,name="rl")
    w=maxw*v/mx
    rect(s,si,x0,yy,w,0.46,col,shape=MSO_SHAPE.ROUNDED_RECTANGLE,name="bar")
    txt(s,si,x0+w+0.12,yy+0.04,1.4,0.4,[[(f"{v:.2f}",15,True,col)]],name="vl")
txt(s,si,4.95,y0+4*rowh-0.02,7.0,0.3,[[("perplexity   (lower = better)",11,True,GRAY)]],name="cap")
rect(s,si,0.40,5.02,6.15,0.92,LG,name="c1"); rect(s,si,0.40,5.02,0.16,0.92,GREEN,name="c1b")
txt(s,si,0.75,5.11,5.7,0.78,[[("−4.2%",24,True,GREEN)],[("vs FLAT-LLM  (faithful angular-BI baseline)",12,False,GRAY)]],name="c1t")
rect(s,si,6.75,5.02,6.15,0.92,LB,name="c2"); rect(s,si,6.75,5.02,0.16,0.92,BLUE,name="c2b")
txt(s,si,7.10,5.11,5.7,0.78,[[("−19.7%",24,True,BLUE)],[("vs Uniform allocation",12,False,GRAY)]],name="c2t")
rect(s,si,0.40,6.10,12.5,0.95,LY,name="why"); rect(s,si,0.40,6.10,0.16,0.95,ORANGE,name="whyb")
txt(s,si,0.75,6.18,12.0,0.82,[[("Why we win — the signals disagree:  ",12,True,ORANGE),
    ("FLAT-LLM's block-influence (how much a layer rotates the residual) favors boundary layers 0 & 31;",12,False,DARK)],
    [("amplification instead protects high-downstream-impact mid-early layers (2–5) and drains the damping tail.  ",12,False,DARK),
    ("Amp captures sensitivity BI misses.",12,True,DARK)]],anchor=MSO_ANCHOR.MIDDLE)

# ===== VALIDATE geometry =====
bad=[]
for (si,name,l,t,w,h) in BOXES:
    if t< -0.001 or l< -0.001 or (t+h)>SH+0.02 or (l+w)>SW+0.02:
        bad.append((si,name,round(l,2),round(t,2),round(l+w,2),round(t+h,2)))
print("OUT-OF-BOUNDS shapes:", len(bad))
for b in bad: print("  ", b)
P.save("hypergraph_slides.pptx")
print("saved; total slides:", len(list(P.slides)))
