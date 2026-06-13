import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle, FancyBboxPatch, Ellipse
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_CONNECTOR

P = Presentation("hypergraph_redundancy.pptx")
SW, SH = 13.3333, 7.5

def rgb(sh):
    try:
        if sh.fill.type is not None and sh.fill.fore_color and sh.fill.fore_color.type is not None:
            c = sh.fill.fore_color.rgb; return (c[0]/255, c[1]/255, c[2]/255)
    except Exception: pass
    return None

def linergb(sh):
    try:
        if sh.line.color and sh.line.color.type is not None:
            c = sh.line.color.rgb; return (c[0]/255, c[1]/255, c[2]/255)
    except Exception: pass
    return None

def trgb(run):
    try:
        if run.font.color and run.font.color.type is not None:
            c = run.font.color.rgb; return (c[0]/255, c[1]/255, c[2]/255)
    except Exception: pass
    return (0, 0, 0)

def shape_name(sh):
    try: return str(sh.auto_shape_type)
    except Exception: return ""

for idx in range(len(P.slides)):
    s = P.slides[idx]
    fig, ax = plt.subplots(1, 1, figsize=(13.33, 7.5))
    ax.set_xlim(0, SW); ax.set_ylim(0, SH); ax.invert_yaxis(); ax.set_aspect('equal')
    ax.add_patch(Rectangle((0, 0), SW, SH, fill=True, facecolor='white', edgecolor='black', lw=2))
    ax.set_xticks([]); ax.set_yticks([])
    for sh in s.shapes:
        # connectors (lines / edges)
        if sh.shape_type == MSO_SHAPE_TYPE.LINE or "Connector" in str(sh.shape_type):
            try:
                x1, y1 = Emu(sh.left).inches, Emu(sh.top).inches
                x2, y2 = x1 + Emu(sh.width).inches, y1 + Emu(sh.height).inches
                lc = linergb(sh) or (0.5, 0.5, 0.5)
                ax.plot([x1, x2], [y1, y2], color=lc, lw=2.2, solid_capstyle='round')
            except Exception: pass
            continue
        l, t = Emu(sh.left).inches, Emu(sh.top).inches
        w, h = Emu(sh.width).inches, Emu(sh.height).inches
        fc = rgb(sh); lc = linergb(sh); nm = shape_name(sh)
        if fc is not None or lc is not None:
            face = fc if fc is not None else 'none'
            edge = lc if lc is not None else 'none'
            if "OVAL" in nm:
                ax.add_patch(Ellipse((l + w/2, t + h/2), w, h, facecolor=face, edgecolor=edge, lw=1.4))
            elif "ROUNDED" in nm:
                ax.add_patch(FancyBboxPatch((l + 0.04, t + 0.04), max(w-0.08, 0.02), max(h-0.08, 0.02),
                             boxstyle="round,pad=0.02,rounding_size=0.08",
                             facecolor=face, edgecolor=edge, lw=1.6))
            else:
                ax.add_patch(Rectangle((l, t), w, h, facecolor=face, edgecolor=edge, lw=1.4))
        if sh.has_text_frame and sh.text_frame.text.strip():
            paras = list(sh.text_frame.paragraphs)
            yi = t + 0.14
            for p in paras:
                if not p.runs:
                    continue
                full = "".join(r.text for r in p.runs)
                r0 = p.runs[0]; sz = (r0.font.size.pt if r0.font.size else 12)
                if not full.strip():
                    yi += sz / 60.0; continue
                al = {1: 'left', 2: 'center', 3: 'right'}.get(p.alignment, 'left') if p.alignment else 'left'
                px = l if al != 'center' else l + w/2
                if al == 'right': px = l + w
                ax.text(px, yi, full[:110], fontsize=min(sz*0.62, 15), color=trgb(r0),
                        ha=al, va='top', fontweight=('bold' if r0.font.bold else 'normal'),
                        style=('italic' if r0.font.italic else 'normal'))
                yi += max(sz / 50.0, 0.2)
    ax.set_title(f"slide {idx+1}", fontsize=9, loc='left')
    plt.tight_layout()
    plt.savefig(f"_slide_{idx+1:02d}.png", dpi=110, bbox_inches='tight')
    plt.close(fig)
print("rendered", len(P.slides), "slides")
