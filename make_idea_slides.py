"""
Idea-only presentation deck for the Typed Pruning Hypergraph.
NO results — pure explanation of the idea, in the visual style of the
Isomorphic Pruning (Fang et al.) figures: clean stacked-box layers, group
brackets, and node-edge graphs.

8 slides:
  1. Title / the setup
  2. Background — three ways to rank what to prune (recreates VainF Fig.2)
  3. Our method in the same 4-column view (parallels slide 2)
  4. Isomorphic pruning in one picture (the baseline we build on)
  5. The gap — what uniform pruning misses (conceptual, no numbers)
  6. VainF's graph vs our graph (Fig.3 style, side by side)
  7. The three parameters, defined (S_min, theta, alpha)
  8. Worked example — building the graph component-by-component
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Colors (from vdms_slides.pptx) ───────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x27, 0x44)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x00, 0x68, 0xB5)
ORANGE  = RGBColor(0xFF, 0xA3, 0x00)
DARK    = RGBColor(0x1A, 0x1A, 0x1A)
BODY    = RGBColor(0x52, 0x52, 0x52)
PANEL   = RGBColor(0xF5, 0xF8, 0xFF)
SKYBLUE = RGBColor(0x9B, 0xBE, 0xDF)
AUTHCLR = RGBColor(0xCC, 0xDD, 0xEE)
PGNUM   = RGBColor(0xAA, 0xAA, 0xAA)
GREEN   = RGBColor(0x70, 0xAD, 0x47)
RED     = RGBColor(0xC0, 0x50, 0x4D)
LGREEN  = RGBColor(0xE2, 0xEF, 0xDA)
LRED    = RGBColor(0xFC, 0xE4, 0xD6)
LBLUE   = RGBColor(0xDE, 0xEB, 0xF7)
LORANGE = RGBColor(0xFF, 0xF2, 0xCC)
GRAY    = RGBColor(0xBF, 0xBF, 0xBF)
LGRAY   = RGBColor(0xED, 0xED, 0xED)
BOXGRAY = RGBColor(0xE6, 0xE6, 0xE6)
EDGEGRAY= RGBColor(0x80, 0x80, 0x80)

F = "Calibri"

SW, SH = 12192000, 6858000
ML, MT = 365760, 228600
CW     = 11430000
RULE_Y = 795528
CONT_Y = 980000
CONT_H = SH - CONT_Y - 228600

# ── primitives ───────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, line_color=None, line_w=12700, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color; s.line.width = Emu(line_w)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def add_tb(slide, x, y, w, h, wrap=True):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tb.text_frame.word_wrap = wrap
    return tb.text_frame

def fp(tf, text, size, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = F; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return p

def ap(tf, text, size, bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=4, italic=False):
    p = tf.add_paragraph(); p.alignment = align
    if space_before: p.space_before = Pt(space_before)
    r = p.add_run(); r.text = text
    r.font.name = F; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return p

def line(slide, x1, y1, x2, y2, color=EDGEGRAY, w=15875, dash=None):
    c = slide.shapes.add_connector(2, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)))
    c.line.color.rgb = color; c.line.width = Emu(w)
    c.shadow.inherit = False
    if dash:
        from pptx.oxml.ns import qn
        ln = c.line._get_or_add_ln()
        d = ln.makeelement(qn('a:prstDash'), {'val': dash})
        ln.append(d)
    return c

_PAGE = [0]   # auto page counter (title slide doesn't call title_rule)

def title_rule(slide, title_text, page_num=None):
    _PAGE[0] += 1
    tf = add_tb(slide, ML, MT, CW, 548640)
    fp(tf, title_text, 28, bold=True, color=DARK)
    add_rect(slide, ML, RULE_Y, CW + 27432, 31750, BLUE)
    tf2 = add_tb(slide, 11430000, 6492240, 640080, 274320)
    fp(tf2, str(_PAGE[0]), 10, color=PGNUM, align=PP_ALIGN.RIGHT)

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def layer_box(slide, x, y, w, h, label=None, fill=BOXGRAY, prune_frac=0.0,
              prune_color=BLUE, label_color=DARK, label_size=10):
    """A pretrained 'layer' box; optionally a colored pruned portion on the right."""
    add_rect(slide, x, y, w, h, fill, line_color=EDGEGRAY, line_w=9525)
    if prune_frac > 0:
        pw = int(w * prune_frac)
        add_rect(slide, x + w - pw, y, pw, h, prune_color, line_color=EDGEGRAY, line_w=9525)
    if label:
        tf = add_tb(slide, x, y + h/2 - 130000, w, 260000)
        fp(tf, label, label_size, color=label_color, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 1 — Title / setup
# =============================================================================

def slide1(prs):
    s = blank(prs)
    add_rect(s, 0, 0, SW, SH, NAVY)

    tf = add_tb(s, ML + 182880, 700000, CW - 365760, 1100000)
    fp(tf, "Typed Pruning Hypergraph", 44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tf2 = add_tb(s, ML + 182880, 1560000, CW - 365760, 600000)
    fp(tf2, "Deciding what to prune by structure, importance, and coupling",
       22, color=SKYBLUE, align=PP_ALIGN.CENTER)

    add_rect(s, 2926080, 2330000, 6309360, 30480, ORANGE)

    box_y, box_h, box_w, gap = 2520000, 760000, 3300000, 200000
    start_x = (SW - 3 * box_w - 2 * gap) // 2
    boxes = [
        ("The model", "A pretrained Vision\nTransformer (DeiT-Small)"),
        ("The task", "Cut its compute by ~45%\nwith minimal accuracy loss"),
        ("The question", "WHICH parts to remove,\nand by HOW MUCH?"),
    ]
    for i, (t, sub) in enumerate(boxes):
        bx = start_x + i * (box_w + gap)
        add_rect(s, bx, box_y, box_w, box_h, BLUE)
        tf_b = add_tb(s, bx + 80000, box_y + 90000, box_w - 160000, box_h - 160000)
        fp(tf_b, t, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, ln_ in enumerate(sub.split("\n")):
            ap(tf_b, ln_, 13, color=SKYBLUE, align=PP_ALIGN.CENTER, space_before=4 if j == 0 else 1)

    tf3 = add_tb(s, ML + 182880, 3560000, CW - 365760, 360000)
    fp(tf3, "Shahrzad Esmat   ·   Iowa State University", 16, color=AUTHCLR, align=PP_ALIGN.CENTER)
    tf4 = add_tb(s, ML + 182880, 3960000, CW - 365760, 360000)
    fp(tf4, "We generalize Isomorphic Pruning (Fang et al.) with a typed dependency graph",
       14, color=ORANGE, align=PP_ALIGN.CENTER, italic=True)


# =============================================================================
# SLIDE 2 — Background: three ways to rank (recreates VainF Fig. 2)
# =============================================================================

def slide2(prs):
    s = blank(prs)
    title_rule(s, "Background: How Do We Decide What to Prune?", 1)

    col_w = 2500000
    gap   = 480000
    start_x = ML + 60000
    top   = CONT_Y + 120000
    bw, bh, bgap = 1700000, 470000, 150000
    nboxes = 4
    box_x_off = (col_w - bw) // 2

    headers = ["(a) Network", "(b) Local Pruning", "(c) Global Pruning", "(d) Isomorphic (grouped)"]
    # per-column box prune fractions + colors
    col_colors = [BLUE, RED, GREEN, ORANGE]

    for c in range(4):
        cx = start_x + c * (col_w + gap)
        # header
        tfh = add_tb(s, cx, top, col_w, 300000)
        fp(tfh, headers[c], 14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        # forward arrow above stack
        ax = cx + col_w // 2
        line(s, ax, top + 330000, ax, top + 430000, color=EDGEGRAY, w=19050)

        boxes_top = top + 460000
        for b in range(nboxes):
            by = boxes_top + b * (bh + bgap)
            bx = cx + box_x_off
            if c == 0:
                layer_box(s, bx, by, bw, bh, fill=BOXGRAY)
            elif c == 1:
                # local: each box ranked independently → its own color slice + tiny bracket
                layer_box(s, bx, by, bw, bh, fill=BOXGRAY,
                          prune_frac=0.30, prune_color=col_colors[b % 4])
                line(s, bx + bw + 40000, by + 40000, bx + bw + 40000, by + bh - 40000,
                     color=col_colors[b % 4], w=15875)
            elif c == 2:
                # global: all boxes one ranking → same blue slice, one big bracket
                layer_box(s, bx, by, bw, bh, fill=BOXGRAY,
                          prune_frac=0.30, prune_color=BLUE)
            elif c == 3:
                # isomorphic: two groups by topology → two colors, two brackets
                grp = BLUE if b in (0, 3) else RED
                layer_box(s, bx, by, bw, bh, fill=BOXGRAY, prune_frac=0.30, prune_color=grp)

        stack_bot = boxes_top + nboxes * (bh + bgap) - bgap
        # brackets / ranking labels
        if c == 1:
            tfl = add_tb(s, cx, stack_bot + 40000, col_w, 260000)
            fp(tfl, "rank within each layer", 11, color=BODY, align=PP_ALIGN.CENTER, italic=True)
        elif c == 2:
            line(s, cx + box_x_off + bw + 60000, boxes_top,
                 cx + box_x_off + bw + 60000, stack_bot, color=BLUE, w=19050)
            tfl = add_tb(s, cx, stack_bot + 40000, col_w, 260000)
            fp(tfl, "one global ranking", 11, color=BLUE, align=PP_ALIGN.CENTER, italic=True)
        elif c == 3:
            # group A bracket (boxes 0..0) and group B; simplify: label
            tfl = add_tb(s, cx, stack_bot + 40000, col_w, 520000)
            fp(tfl, "group by structure type,", 11, color=DARK, align=PP_ALIGN.CENTER, italic=True)
            ap(tfl, "rank within each group", 11, color=ORANGE, align=PP_ALIGN.CENTER, italic=True, space_before=1)

    # caption
    cap = add_tb(s, ML, SH - 760000, CW, 520000)
    fp(cap, "Comparing importance across different structure types is unfair (different scales). "
            "Isomorphic pruning groups parameters by structural type and ranks WITHIN each group.",
       13, color=BODY, italic=True)


# =============================================================================
# SLIDE 3 — Isomorphic pruning in one picture
# =============================================================================

def slide3(prs):
    s = blank(prs)
    title_rule(s, "Isomorphic Pruning: One Ratio Per Structure Type", 3)

    # left panel text
    add_rect(s, 274320, CONT_Y, 5300000, CONT_H, PANEL)
    add_rect(s, 274320, CONT_Y, 76200, CONT_H, BLUE)
    tx = 457200 + 60000
    tf = add_tb(s, tx, CONT_Y + 80000, 4800000, CONT_H - 160000)
    fp(tf, "The baseline we build on", 17, bold=True, color=BLUE)
    ap(tf, "", 6)
    ap(tf, "Two structure types, each pruned by ONE ratio:", 14, bold=True, color=DARK, space_before=4)
    ap(tf, "  • Attention — shrink head_dim inside each head", 13, color=BODY, space_before=3)
    ap(tf, "  • MLP — shrink the hidden dimension", 13, color=BODY, space_before=2)
    ap(tf, "", 5)
    ap(tf, "Every one of the 12 blocks gets the SAME ratio.", 14, bold=True, color=DARK, space_before=4)
    ap(tf, "", 4)
    ap(tf, "The residual stream (embed dim) is never cut —", 13, color=BODY, space_before=2)
    ap(tf, "only the internals of each block shrink.", 13, color=BODY, space_before=1)
    ap(tf, "", 6)
    ap(tf, "Simple and effective — but it treats a critical", 13, color=DARK, space_before=4, italic=True)
    ap(tf, "block exactly like a redundant one.", 13, color=DARK, italic=True, space_before=1)

    # right: 12 identical blocks, each split [Attn | MLP], same prune everywhere
    rx = 6100000
    ry = CONT_Y + 120000
    tfh = add_tb(s, rx, ry - 40000, 5600000, 260000)
    fp(tfh, "All 12 blocks — identical treatment", 14, bold=True, color=DARK)

    bw, bh, bgap = 4600000, 340000, 70000
    ry += 280000
    attn_w = int(bw * 0.42)
    for i in range(12):
        by = ry + i * (bh + bgap)
        # attention half
        layer_box(s, rx, by, attn_w, bh, fill=LBLUE)
        add_rect(s, rx + attn_w - int(attn_w*0.30), by, int(attn_w*0.30), bh, BLUE,
                 line_color=EDGEGRAY, line_w=9525)
        # mlp half
        layer_box(s, rx + attn_w + 40000, by, bw - attn_w - 40000, bh, fill=LORANGE)
        mw = bw - attn_w - 40000
        add_rect(s, rx + attn_w + 40000 + mw - int(mw*0.30), by, int(mw*0.30), bh, ORANGE,
                 line_color=EDGEGRAY, line_w=9525)
        # labels
        tfb = add_tb(s, rx + 60000, by + bh/2 - 120000, attn_w, 240000)
        fp(tfb, f"Block {i}  ·  Attn", 10, color=DARK)
        tfm = add_tb(s, rx + attn_w + 100000, by + bh/2 - 120000, mw, 240000)
        fp(tfm, "MLP", 10, color=DARK)

    tfn = add_tb(s, rx, ry + 12 * (bh + bgap) + 20000, 5600000, 260000)
    fp(tfn, "blue / orange = pruned portion — same fraction in every block",
       11, color=BODY, italic=True)


# =============================================================================
# SLIDE 4 — The gap (conceptual)
# =============================================================================

def slide4(prs):
    s = blank(prs)
    title_rule(s, "What Uniform Pruning Misses")

    sub = add_tb(s, ML, CONT_Y, CW, 320000)
    fp(sub, "One ratio for every block cannot express three real differences between blocks:",
       15, bold=True, color=DARK)

    # (num, title, accent, fill, [detail lines], gap_metric, fix_symbol, fix_text)
    cards = [
        ("1", "Some blocks are nearly dead weight", RED, LRED,
         ["Block sensitivity (output change when a block is bypassed) varies widely —",
          "a few blocks are close to identity, yet uniform pruning still keeps and shrinks them,",
          "spending MAC budget on blocks that contribute almost nothing."],
         "depth ignored", "S_min", "remove them entirely"),
        ("2", "Blocks are not equally important", ORANGE, LORANGE,
         ["Taylor importance (|grad × weight|) differs block-to-block. One global ratio",
          "over-prunes the critical blocks (accuracy drops) and under-prunes the redundant",
          "ones (budget wasted) — the cut lands in exactly the wrong places."],
         "width is flat", "θ", "per-group ratios"),
        ("3", "Block importances are not independent", BLUE, LBLUE,
         ["Some blocks rise and fall in importance together (their scores are correlated).",
          "Pruning each block in isolation ignores that a surviving block may depend on a",
          "neighbour you just pruned — the decisions should inform each other."],
         "coupling ignored", "α", "functional edges"),
    ]
    cw = CW
    ch = 1430000
    cgap = 150000
    cy = CONT_Y + 380000
    chip_w = 2550000
    for i, (num, head, accent, fill, lines, metric, sym, fix) in enumerate(cards):
        y = cy + i * (ch + cgap)
        add_rect(s, ML, y, cw, ch, fill)
        add_rect(s, ML, y, 150000, ch, accent)
        # number circle
        add_rect(s, ML + 300000, y + ch/2 - 340000, 680000, 680000, accent, shape=MSO_SHAPE.OVAL)
        tfn = add_tb(s, ML + 300000, y + ch/2 - 280000, 680000, 560000)
        fp(tfn, num, 32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # title + detail
        tf = add_tb(s, ML + 1180000, y + 110000, cw - 1300000 - chip_w, ch - 220000)
        fp(tf, head, 18, bold=True, color=accent)
        fp_first = True
        for ln_ in lines:
            ap(tf, ln_, 13, color=BODY, space_before=6 if fp_first else 2)
            fp_first = False
        # gap metric tag (italic, under title region)
        tfm = add_tb(s, ML + 1180000, y + ch - 300000, 3000000, 240000)
        fp(tfm, f"isomorphic: {metric}", 12, italic=True, color=accent)
        # right "fix" chip
        chx = ML + cw - chip_w - 120000
        add_rect(s, chx, y + 150000, chip_w, ch - 300000, WHITE, line_color=accent, line_w=19050)
        tfc = add_tb(s, chx + 120000, y + 230000, chip_w - 240000, ch - 420000)
        fp(tfc, "our fix", 11, bold=True, color=BODY)
        ap(tfc, f"→  {sym}", 24, bold=True, color=accent, space_before=4)
        ap(tfc, fix, 13, color=DARK, space_before=4)


# =============================================================================
# SLIDE 5 — Our idea: Typed Pruning Hypergraph (Fig.3 style graph)
# =============================================================================

def _func_arcs(s, cx, node_d, base_y, functional, w, hi, step, label=False):
    """Draw dashed functional-edge arcs, STAGGERED by span so they don't merge:
    the widest arc sits highest (outer) and narrower arcs nest inside it."""
    order = sorted(functional, key=lambda e: -(e[1] - e[0]))   # widest first
    for rank, (i, j) in enumerate(order):
        x1 = cx[i] + node_d / 2
        x2 = cx[j] + node_d / 2
        ytop = base_y - (hi - rank * step)
        for (a, b, c, d2) in [(x1, base_y, x1, ytop), (x1, ytop, x2, ytop), (x2, ytop, x2, base_y)]:
            line(s, a, b, c, d2, color=ORANGE, w=w, dash='dash')
        if label:
            tfm = add_tb(s, (x1 + x2) // 2 - 450000, ytop - 220000, 900000, 200000)
            fp(tfm, "w_ij", 11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, italic=True)


def _graph_panel(s, px, pw, title, title_fill, removed, functional, node_fills, node_top):
    """Draw one block-chain graph panel. removed=set of indices, functional=list
    of (i,j) dashed coupling arcs, node_fills=list of per-node fill colors."""
    add_rect(s, px, CONT_Y + 60000, pw, 380000, title_fill)
    tfh = add_tb(s, px, CONT_Y + 100000, pw, 300000)
    fp(tfh, title, 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    n = 4
    node_d = 820000
    inner = pw - 360000
    node_gap = (inner - n * node_d) // (n - 1)
    x0 = px + 180000
    ny = node_top
    cx = [x0 + i * (node_d + node_gap) for i in range(n)]

    # residual chain (solid gray) between consecutive nodes
    for i in range(n - 1):
        line(s, cx[i] + node_d, ny + node_d/2, cx[i+1], ny + node_d/2,
             color=EDGEGRAY, w=22225)

    # functional dashed arcs (above) — our addition; staggered so they don't merge
    _func_arcs(s, cx, node_d, ny, functional, w=20320, hi=470000, step=210000, label=False)

    # nodes (clean circles — E_s is explained in the legend)
    for i in range(n):
        if i in removed:
            add_rect(s, cx[i], ny, node_d, node_d, LGRAY, line_color=GRAY,
                     line_w=12700, shape=MSO_SHAPE.OVAL)
            tfx = add_tb(s, cx[i], ny + node_d/2 - 220000, node_d, 440000)
            fp(tfx, "✕", 32, bold=True, color=RED, align=PP_ALIGN.CENTER)
            tfr = add_tb(s, cx[i] - 60000, ny + node_d + 40000, node_d + 120000, 240000)
            fp(tfr, "removed", 11, color=RED, align=PP_ALIGN.CENTER, italic=True)
        else:
            add_rect(s, cx[i], ny, node_d, node_d, node_fills[i], shape=MSO_SHAPE.OVAL)
            tfn = add_tb(s, cx[i], ny + node_d/2 - 160000, node_d, 320000)
            fp(tfn, f"Blk {i}", 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def _legend_row(s, x, y, kind, term, defn, term_color):
    """One legend row: visual swatch + bold term + definition."""
    sw_w = 520000
    if kind == "solid":
        line(s, x, y + 90000, x + sw_w, y + 90000, color=EDGEGRAY, w=22225)
    elif kind == "struct":
        for k, cc in enumerate([SKYBLUE, BLUE, NAVY]):
            add_rect(s, x + k*170000, y + 10000, 140000, 140000, cc,
                     line_color=EDGEGRAY, line_w=6350)
    elif kind == "dash":
        line(s, x, y + 90000, x + sw_w, y + 90000, color=ORANGE, w=19050, dash='dash')
    tf = add_tb(s, x + sw_w + 140000, y - 40000, CW - sw_w - 700000, 360000)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = term + "  —  "
    r.font.name = F; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = term_color
    r2 = p.add_run(); r2.text = defn
    r2.font.name = F; r2.font.size = Pt(13); r2.font.color.rgb = BODY


def slide5(prs):
    s = blank(prs)
    title_rule(s, "The Key Difference: VainF's Graph  vs  Our Graph", 5)

    half = SW // 2
    lpx, lpw = ML, half - ML - 120000
    rpx, rpw = half + 120000, SW - (half + 120000) - ML
    node_top = CONT_Y + 1000000

    # LEFT — VainF: structural edges only, no removal, uniform color
    _graph_panel(s, lpx, lpw, "Isomorphic Pruning (VainF)", BLUE,
                 removed=set(), functional=[], node_fills=[BLUE]*4, node_top=node_top)
    cl = add_tb(s, lpx, node_top + 1180000, lpw, 360000)
    fp(cl, "Blocks pruned in isolation · one ratio per type · all kept",
       12, color=BODY, align=PP_ALIGN.CENTER, italic=True)

    # RIGHT — Ours
    _graph_panel(s, rpx, rpw, "Our Typed Hypergraph  H = (V', E_s, E_f)", ORANGE,
                 removed={2}, functional=[(0, 3), (1, 3)],
                 node_fills=[GREEN, GREEN, None, RED], node_top=node_top)
    cr = add_tb(s, rpx, node_top + 1180000, rpw, 360000)
    fp(cr, "S_min removes a block · θ colors groups · α adds w_ij edges",
       12, color=BODY, align=PP_ALIGN.CENTER, italic=True)

    # ── clean edge legend (defines each line/edge precisely) ──────────────────
    leg_y = node_top + 1640000
    add_rect(s, ML, leg_y, CW, 1180000, PANEL)
    add_rect(s, ML, leg_y, 76200, 1180000, BLUE)
    th = add_tb(s, ML + 200000, leg_y + 40000, CW - 400000, 260000)
    fp(th, "How to read the graph", 14, bold=True, color=BLUE)
    rx = ML + 240000
    _legend_row(s, rx, leg_y + 350000, "solid", "Residual stream",
                "data flowing from one block to the next  (both methods)", DARK)
    _legend_row(s, rx, leg_y + 620000, "struct", "E_s  Structural edge",
                "weights INSIDE one block that must be cut together (qkv·proj·mlp)  —  both methods", BLUE)
    _legend_row(s, rx, leg_y + 890000, "dash", "E_f  Functional edge",
                "links TWO blocks with similar importance; lets importance flow between them  —  NEW (ours)", ORANGE)

    # bottom takeaway
    by = SH - 470000
    add_rect(s, ML, by, CW, 330000, NAVY)
    add_rect(s, ML, by, 120000, 330000, ORANGE)
    tf = add_tb(s, ML + 280000, by + 40000, CW - 450000, 260000)
    fp(tf, "Same structural edges as VainF  +  block removal (V')  +  functional edges (E_f).",
       15, bold=True, color=WHITE)


# =============================================================================
# SLIDE 6 — Three knobs + punchline
# =============================================================================

def slide6(prs):
    s = blank(prs)
    title_rule(s, "The Three Parameters, Defined", 6)

    # (symbol, name, accent, fill, light_formula_fill, definition, [formula lines], effect)
    cards = [
        ("S_min", "depth-pruning threshold", RED, LRED, RGBColor(0xF7,0xD4,0xCC),
         "A threshold on block sensitivity — how much the output changes when a block is skipped.",
         ["S(i) = mean ‖f(x) − f_bypass(x)‖ / ‖f(x)‖",
          "remove block i  if  S(i) < S_min"],
         "Graph effect:  shrinks the node set V′ — redundant blocks are deleted entirely."),
        ("θ", "grouping threshold (theta)", BLUE, LBLUE, RGBColor(0xC9,0xDF,0xF2),
         "Blocks with similar importance are merged into one group; each group gets its own ratio.",
         ["group i, j   if   |Î_i − Î_j| < θ",
          "θ = 1 → one group (uniform);  θ = 0 → all separate"],
         "Graph effect:  partitions blocks into groups — important groups are pruned less."),
        ("α", "coupling strength (alpha)", ORANGE, LORANGE, RGBColor(0xFF,0xE8,0xB8),
         "Adds functional edges between blocks of similar importance and lets them boost each other.",
         ["w_ij = min(I_i, I_j) / max(I_i, I_j)",
          "I↑(j) = I(j) · (1 + α · Σ w_ij·I(i)/Z)"],
         "Graph effect:  creates E_f — a block's importance rises with its coupled neighbours'."),
    ]
    cw = 3700000
    cgap = 180000
    start_x = (SW - 3 * cw - 2 * cgap) // 2
    cy = CONT_Y + 40000
    ch = 3680000
    for i, (sym, sub, accent, fill, fbox, defn, formula, effect) in enumerate(cards):
        x = start_x + i * (cw + cgap)
        add_rect(s, x, cy, cw, ch, fill)
        add_rect(s, x, cy, cw, 140000, accent)

        # symbol + name
        th = add_tb(s, x + 200000, cy + 220000, cw - 360000, 720000)
        fp(th, sym, 30, bold=True, color=accent)
        ap(th, sub, 13, bold=True, color=DARK, space_before=2)

        # definition
        td = add_tb(s, x + 200000, cy + 960000, cw - 380000, 820000)
        fp(td, "Definition", 11, bold=True, color=accent)
        ap(td, defn, 12.5, color=BODY, space_before=3)

        # formula box
        fby = cy + 1880000
        add_rect(s, x + 160000, fby, cw - 320000, 620000, fbox)
        tff = add_tb(s, x + 240000, fby + 60000, cw - 480000, 500000)
        fp(tff, formula[0], 12, bold=True, color=accent, italic=True)
        if len(formula) > 1:
            ap(tff, formula[1], 12, bold=True, color=accent, italic=True, space_before=6)

        # effect
        te = add_tb(s, x + 200000, cy + 2620000, cw - 380000, 980000)
        fp(te, effect, 12.5, color=DARK)

    # punchline box
    py = cy + ch + 180000
    add_rect(s, ML, py, CW, 680000, NAVY)
    add_rect(s, ML, py, 150000, 680000, ORANGE)
    tf = add_tb(s, ML + 320000, py + 70000, CW - 500000, 560000)
    fp(tf, "Set  S_min = 0,  θ = 1,  α = 0  →  exactly Isomorphic Pruning (VainF).",
       17, bold=True, color=WHITE)
    ap(tf, "VainF is one corner of our space; the three knobs let us search for a better point.",
       15, color=ORANGE, space_before=6)


def _mini_graph(s, gx, gy, node_d, gap, fills, removed, functional, labels):
    """Compact 5-node block-chain for the build-up example."""
    n = len(fills)
    cx = [gx + i * (node_d + gap) for i in range(n)]
    present = [i for i in range(n) if i not in removed]

    # residual chain between consecutive present nodes (skips removed)
    for a, b in zip(present, present[1:]):
        line(s, cx[a] + node_d, gy + node_d/2, cx[b], gy + node_d/2,
             color=EDGEGRAY, w=17780)

    # functional dashed arcs above — staggered so distinct couplings don't merge
    _func_arcs(s, cx, node_d, gy, functional, w=15875, hi=300000, step=150000, label=False)

    # nodes
    for i in range(n):
        if i in removed:
            add_rect(s, cx[i], gy, node_d, node_d, LGRAY, line_color=GRAY,
                     line_w=9525, shape=MSO_SHAPE.OVAL)
            tfx = add_tb(s, cx[i], gy + node_d/2 - 150000, node_d, 300000)
            fp(tfx, "✕", 20, bold=True, color=RED, align=PP_ALIGN.CENTER)
        else:
            add_rect(s, cx[i], gy, node_d, node_d, fills[i], shape=MSO_SHAPE.OVAL)
            tfn = add_tb(s, cx[i], gy + node_d/2 - 120000, node_d, 240000)
            fp(tfn, labels[i], 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide7(prs):
    s = blank(prs)
    title_rule(s, "Worked Example: Building the Graph, One Component at a Time", 7)

    labels = ["B0", "B1", "B2", "B3", "B4"]
    # illustrative values (to show the mechanism, not measured results)
    # sensitivity S: B0 .90  B1 .50  B2 .35  B3 .60  B4 .85   (S_min = 0.40 removes B2)
    # importance  I: B0 .45  B1 .22  B3 .25  B4 .40           (groups: {B0,B4} hi, {B1,B3} lo)

    rows = [
        ("0   Start  =  VainF",
         "All 5 blocks kept · one uniform ratio · no inter-block edges",
         [BLUE]*5, set(), []),
        ("+  S_min = 0.40",
         "B2 has S = 0.35 < 0.40  →  removed entirely  (node set V′ shrinks)",
         [BLUE, BLUE, BLUE, BLUE, BLUE], {2}, []),
        ("+  θ  (grouping)",
         "Group by importance: {B0,B4} high → pruned less,  {B1,B3} low → pruned more",
         [GREEN, RED, BLUE, RED, GREEN], {2}, []),
        ("+  α  (coupling)",
         "Functional edges link similar blocks: B0–B4 and B1–B3 boost each other",
         [GREEN, RED, BLUE, RED, GREEN], {2}, [(0, 4), (1, 3)]),
    ]

    ry0 = CONT_Y + 40000
    row_h = 1180000
    lab_w = 3250000
    gx = ML + lab_w + 250000
    node_d = 470000
    graph_w = SW - gx - ML
    gap = (graph_w - 5 * node_d) // 4

    for r, (head, change, fills, removed, functional) in enumerate(rows):
        ry = ry0 + r * row_h
        # alternating light background band
        if r % 2 == 0:
            add_rect(s, ML, ry, CW, row_h - 60000, RGBColor(0xF7, 0xF9, 0xFC))
        # left label
        accent = [GRAY, RED, BLUE, ORANGE][r]
        add_rect(s, ML, ry + 120000, 110000, row_h - 320000, accent)
        tf = add_tb(s, ML + 200000, ry + 130000, lab_w - 250000, row_h - 260000)
        fp(tf, head, 16, bold=True, color=accent)
        ap(tf, change, 12, color=BODY, space_before=6)
        # graph
        gy = ry + (row_h - node_d) // 2 + 40000
        _mini_graph(s, gx, gy, node_d, gap, fills, removed, functional, labels)

    # footnote
    fn = add_tb(s, ML, SH - 470000, CW, 320000)
    fp(fn, "Illustrative 5-block example (values chosen to show the mechanism). "
           "Each row adds one component to the row above — the graph is built up, not rebuilt.",
       12, color=BODY, italic=True)


def slide_ours_view(prs):
    """Same 4-column view as slide 2, but for OUR method's progression."""
    s = blank(prs)
    title_rule(s, "Our Method in the Same View: Extending Isomorphic")

    col_w, gap = 2500000, 480000
    start_x = ML + 60000
    top = CONT_Y + 120000
    bw, bh, bgap = 1700000, 460000, 165000
    nboxes = 4
    box_x_off = (col_w - bw) // 2

    headers = ["(a) Pretrained", "(b) + S_min  (depth)", "(c) + θ  (groups)", "(d) + α  (coupling)"]
    removed_box = 2
    group_green = {0, 1}          # high-importance group → pruned LESS
    group_red   = {3}             # low-importance group  → pruned MORE

    for c in range(4):
        cx = start_x + c * (col_w + gap)
        tfh = add_tb(s, cx, top, col_w, 300000)
        fp(tfh, headers[c], 14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        ax = cx + col_w // 2
        line(s, ax, top + 330000, ax, top + 430000, color=EDGEGRAY, w=19050)

        boxes_top = top + 460000
        mids = []
        for b in range(nboxes):
            by = boxes_top + b * (bh + bgap)
            bx = cx + box_x_off
            mids.append(by + bh // 2)
            removed = (c >= 1 and b == removed_box)
            if removed:
                add_rect(s, bx, by, bw, bh, LGRAY, line_color=GRAY, line_w=9525)
                tfx = add_tb(s, bx, by + bh/2 - 140000, bw, 280000)
                fp(tfx, "✕  removed", 12, bold=True, color=RED, align=PP_ALIGN.CENTER)
                continue
            if c == 0:
                layer_box(s, bx, by, bw, bh, fill=BOXGRAY)
                tag = None
            elif c == 1:
                layer_box(s, bx, by, bw, bh, fill=BOXGRAY, prune_frac=0.30, prune_color=BLUE)
                tag = None
            else:
                if b in group_green:
                    layer_box(s, bx, by, bw, bh, fill=LGREEN, prune_frac=0.18, prune_color=GREEN)
                    tag = ("r ↓", GREEN)
                else:
                    layer_box(s, bx, by, bw, bh, fill=LRED, prune_frac=0.48, prune_color=RED)
                    tag = ("r ↑", RED)
            # block label (left) + optional ratio tag
            tfl = add_tb(s, bx + 40000, by + bh/2 - 130000, 700000, 260000)
            fp(tfl, f"B{b}", 12, bold=True, color=DARK)
            if tag:
                tft = add_tb(s, bx + bw - 560000, by + bh/2 - 130000, 520000, 260000)
                fp(tft, tag[0], 12, bold=True, color=tag[1], align=PP_ALIGN.RIGHT)

        stack_bot = boxes_top + nboxes * (bh + bgap) - bgap
        # functional coupling arc in (d): link the two adjacent green-group boxes
        if c == 3:
            bxr = cx + box_x_off + bw
            y0, y1 = mids[0], mids[1]
            ox = bxr + 120000
            for seg in [(bxr, y0, ox, y0), (ox, y0, ox, y1), (ox, y1, bxr, y1)]:
                line(s, *seg, color=ORANGE, w=17780, dash='dash')
            tfw = add_tb(s, ox + 30000, (y0 + y1)//2 - 110000, 520000, 220000)
            fp(tfw, "w_ij", 10, bold=True, color=ORANGE, italic=True)

        notes = ["all blocks kept,\none uniform ratio",
                 "B2 removed entirely\n(S < S_min)",
                 "groups → own ratio\ngreen r↓ , red r↑",
                 "couple same group\nB0–B1  (edge E_f)"]
        tfn = add_tb(s, cx, stack_bot + 60000, col_w, 520000)
        fp(tfn, notes[c].split("\n")[0], 11, color=BODY, align=PP_ALIGN.CENTER, italic=True)
        ap(tfn, notes[c].split("\n")[1], 11, color=BODY, align=PP_ALIGN.CENTER, italic=True, space_before=1)

    cap = add_tb(s, ML, SH - 640000, CW, 420000)
    fp(cap, "Same column view as the field's — but each step here is OUR addition. "
            "Turn all three off (no removal, one group, no edges) → back to isomorphic pruning.",
       13, color=BODY, italic=True)


def slide_ours_picture(prs):
    """Parallel to the isomorphic 'one ratio per type' slide, but for OUR method:
    heterogeneous — depth removal + per-group width + coupling, residual untouched."""
    s = blank(prs)
    title_rule(s, "Our Pruning: Depth + Per-Group Width + Coupling")

    # left panel text
    add_rect(s, 274320, CONT_Y, 5300000, CONT_H, PANEL)
    add_rect(s, 274320, CONT_Y, 76200, CONT_H, ORANGE)
    tx = 457200 + 60000
    tf = add_tb(s, tx, CONT_Y + 80000, 4800000, CONT_H - 160000)
    fp(tf, "Heterogeneous by design", 17, bold=True, color=ORANGE)
    ap(tf, "", 6)
    ap(tf, "Same two structure types (Attn, MLP), but the", 13, color=DARK, space_before=4)
    ap(tf, "ratio now varies block-to-block:", 13, color=DARK, space_before=1)
    ap(tf, "", 4)
    ap(tf, "• S_min — remove low-sensitivity blocks", 13, bold=True, color=RED, space_before=3)
    ap(tf, "    (here blocks 3, 4, 5 are deleted)", 12, color=BODY, space_before=1)
    ap(tf, "• θ — each importance group its own ratio", 13, bold=True, color=GREEN, space_before=3)
    ap(tf, "    green = important → pruned less (thin slice)", 12, color=BODY, space_before=1)
    ap(tf, "    red = redundant → pruned more (thick slice)", 12, color=BODY, space_before=1)
    ap(tf, "• α — functional edges couple similar blocks", 13, bold=True, color=ORANGE, space_before=3)
    ap(tf, "", 5)
    ap(tf, "Residual stream (embed dim) still never cut —", 13, color=DARK, space_before=3, italic=True)
    ap(tf, "only block internals shrink, just like isomorphic.", 13, color=DARK, italic=True, space_before=1)
    ap(tf, "", 5)
    ap(tf, "Same MAC budget — spent where it costs least.", 14, bold=True, color=ORANGE, space_before=4)

    # right: 12 blocks, heterogeneous (depth removals + per-group slices)
    rx = 6100000
    ry = CONT_Y + 120000
    tfh = add_tb(s, rx, ry - 40000, 5600000, 260000)
    fp(tfh, "All 12 blocks — treated by importance", 14, bold=True, color=DARK)

    bw, bh, bgap = 4250000, 300000, 56000
    ry += 250000
    attn_w = int(bw * 0.42)
    removed = {3, 4, 5}
    green   = {0, 1, 2, 11}
    # per-group (attn_frac, mlp_frac)
    frac = {"g": (0.12, 0.18), "r": (0.30, 0.55)}
    mids = {}

    for i in range(12):
        by = ry + i * (bh + bgap)
        mids[i] = by + bh // 2
        if i in removed:
            add_rect(s, rx, by, bw, bh, LGRAY, line_color=GRAY, line_w=9525)
            tfx = add_tb(s, rx + 200000, by + bh/2 - 120000, bw, 240000)
            fp(tfx, f"Block {i}   ✕ removed  (S < S_min)", 10, bold=True, color=RED)
            continue
        g = "g" if i in green else "r"
        fa, fm = frac[g]
        gcol = GREEN if g == "g" else RED
        # group tag bar on the far left
        add_rect(s, rx - 90000, by, 60000, bh, gcol)
        # attention half
        layer_box(s, rx, by, attn_w, bh, fill=LBLUE)
        add_rect(s, rx + attn_w - int(attn_w*fa), by, int(attn_w*fa), bh, BLUE,
                 line_color=EDGEGRAY, line_w=9525)
        # mlp half
        mw = bw - attn_w - 40000
        layer_box(s, rx + attn_w + 40000, by, mw, bh, fill=LORANGE)
        add_rect(s, rx + attn_w + 40000 + mw - int(mw*fm), by, int(mw*fm), bh, ORANGE,
                 line_color=EDGEGRAY, line_w=9525)
        # labels
        tfb = add_tb(s, rx + 50000, by + bh/2 - 110000, attn_w, 220000)
        fp(tfb, f"Blk {i} · Attn", 9, color=DARK)
        tfm = add_tb(s, rx + attn_w + 90000, by + bh/2 - 110000, mw, 220000)
        fp(tfm, "MLP", 9, color=DARK)

    # one functional coupling arc between two adjacent green blocks (0,1)
    bxr = rx + bw
    ox = bxr + 130000
    for seg in [(bxr, mids[0], ox, mids[0]), (ox, mids[0], ox, mids[1]), (ox, mids[1], bxr, mids[1])]:
        line(s, *seg, color=ORANGE, w=17780, dash='dash')
    tfw = add_tb(s, ox + 30000, (mids[0]+mids[1])//2 - 110000, 520000, 220000)
    fp(tfw, "w_ij", 10, bold=True, color=ORANGE, italic=True)

    tfn = add_tb(s, rx, ry + 12 * (bh + bgap) + 20000, 5800000, 260000)
    fp(tfn, "thin slice = pruned less (green) · thick slice = pruned more (red) · grey = removed",
       11, color=BODY, italic=True)


def main():
    prs = Presentation()
    prs.slide_width = Emu(SW); prs.slide_height = Emu(SH)
    _PAGE[0] = 0
    slide1(prs); slide2(prs); slide_ours_view(prs)
    slide3(prs); slide_ours_picture(prs)
    slide4(prs); slide5(prs); slide6(prs); slide7(prs)
    out = "/work/hdd/bdjd/hypergraph_pruning/hypergraph_slides.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
