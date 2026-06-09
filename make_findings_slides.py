from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

P = Presentation("hypergraph_slides.pptx")
BLANK = P.slide_layouts[6]
DARK=RGBColor(0x1A,0x1A,0x1A); GRAY=RGBColor(0x52,0x52,0x52); MUT=RGBColor(0xAA,0xAA,0xAA)
WHITE=RGBColor(0xFF,0xFF,0xFF); BLUE=RGBColor(0x00,0x68,0xB5); RED=RGBColor(0xC0,0x50,0x4D)
ORANGE=RGBColor(0xFF,0xA3,0x00); GREEN=RGBColor(0x2E,0x8B,0x57)
LB=RGBColor(0xDE,0xEB,0xF7); LY=RGBColor(0xFF,0xF2,0xCC); LP=RGBColor(0xFC,0xE4,0xD6)
LG=RGBColor(0xE2,0xF0,0xD9)

def rect(s,l,t,w,h,fill,shape=MSO_SHAPE.RECTANGLE,line=None):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False
    return sp

def txt(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True,space=None):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align
        if space is not None: p.space_after=Pt(space)
        for (text,sz,bold,col) in line:
            r=p.add_run(); r.text=text; r.font.size=Pt(sz); r.font.bold=bold
            r.font.name="Calibri"; r.font.color.rgb=col
    return tb

def header(s,title,num):
    txt(s,0.40,0.25,12.5,0.6,[[(title,28,True,DARK)]])
    rect(s,0.40,0.87,12.53,0.03,BLUE)
    txt(s,12.50,7.10,0.70,0.30,[[(str(num),10,False,MUT)]],align=PP_ALIGN.RIGHT)

def base_slide():
    return P.slides.add_slide(BLANK)

# ============ SLIDE 10 : FINDING 1 — amplification curve ============
s=base_slide(); header(s,"Finding 1:  Compression Error Amplifies with Depth",10)
txt(s,0.40,1.02,12.5,0.4,[[("A layer's ",15,True,DARK),("local",15,True,RED),
    (" error is not its true cost — it propagates to the output, and the amplification depends on depth.",15,True,DARK)]])
# method strip
rect(s,0.40,1.55,12.5,0.62,LB)
txt(s,0.62,1.62,12.1,0.5,[[("Probe:  ",12,True,BLUE),
    ("rank-truncate ONE layer's FFN, measure ",12,False,GRAY),
    ("Δ end-to-end output ÷ Δ local output",12,True,DARK),
    ("   ( = how much that layer's error blows up by the time it reaches the logits )",12,False,GRAY)]])
# bar chart of amplification (error-prop probe, FRAC=0.30, Llama-2-7B)
data=[("L0",30.8),("L1",14.9),("L3",13.4),("L5",5.4),("L10",2.0),("L16",1.8),("L24",1.2),("L31",0.7)]
x0,y0,bw,gap,maxh,base=0.95,5.95,1.05,0.40,3.1,2.45
mx=math.log10(max(v for _,v in data)+1)
for i,(lab,v) in enumerate(data):
    h=maxh*math.log10(v+1)/mx
    cx=x0+i*(bw+gap)
    col=RED if v>=5 else (ORANGE if v>=1.3 else BLUE)
    rect(s,cx,base-h,bw,h,col,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s,cx-0.15,base-h-0.32,bw+0.3,0.3,[[(f"{v:.1f}×",13,True,col)]],align=PP_ALIGN.CENTER)
    txt(s,cx-0.15,base+0.04,bw+0.3,0.25,[[(lab,11,False,GRAY)]],align=PP_ALIGN.CENTER)
txt(s,0.62,2.55,4.0,0.25,[[("amplification (Δe2e / Δlocal)",11,True,GRAY)]])
txt(s,0.62,5.65,12.0,0.25,[[("early layers  →  amplify ~30×",11,True,RED),
    ("        late layers  →  damp (<1×)",11,True,BLUE)]])
# takeaway
rect(s,0.40,6.45,12.5,0.72,LP); rect(s,0.40,6.45,0.16,0.72,RED)
txt(s,0.75,6.52,12.0,0.6,[[("Implication:  ",13,True,RED),
    ("the true cost of compressing a layer varies ~43× by depth. Methods that allocate rank by ",13,False,DARK),
    ("local",13,True,DARK),(" error (FLAT-LLM, ASVD, SVD-LLM) optimize the wrong quantity.",13,False,DARK)]])

# ============ SLIDE 11 : FINDING 2 — banded coupling graph ============
s=base_slide(); header(s,"Finding 2:  Layer Errors Couple — a Graph, Not Just a Curve",11)
txt(s,0.40,1.02,12.5,0.4,[[("Do layers' compression errors ",15,True,DARK),("interact",15,True,BLUE),
    ("?  If so, rank must be allocated jointly, not greedily per-layer.",15,True,DARK)]])
cards=[
 ("Errors are non-additive",ORANGE,LY,"ι",
  "Compressing two layers together ≠ sum of each alone.  Interaction ι (mean 0.60) is large — joint effects matter."),
 ("Coupling is local (banded)",BLUE,LB,"|cos|",
  "Error directions align for nearby layers, decay with distance:  adjacent 0.46  ›  distant 0.28.  A banded graph."),
 ("The tail decouples",RED,LP,"28–31",
  "Final block (layers 28–31) is weakly coupled to all others (|cos| 0.12–0.24) — a separate, compressible group."),
]
y=1.62
for i,(h,acc,fill,big,body) in enumerate(cards):
    rect(s,0.40,y,12.5,1.36,fill); rect(s,0.40,y,0.16,1.36,acc)
    rect(s,0.73,y+0.31,0.74,0.74,acc,shape=MSO_SHAPE.OVAL)
    txt(s,0.73,y+0.34,0.74,0.7,[[(big,18 if len(big)>2 else 26,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,1.69,y+0.18,10.9,0.4,[[(h,18,True,acc)]])
    txt(s,1.69,y+0.62,10.9,0.7,[[(body,13,False,GRAY)]])
    y+=1.50
rect(s,0.40,6.45,12.5,0.72,LB); rect(s,0.40,6.45,0.16,0.72,BLUE)
txt(s,0.75,6.52,12.0,0.6,[[("The adjacency is real:  ",13,True,BLUE),
    ("a banded error-coupling graph over the residual stream.  Rank allocation is a ",13,False,DARK),
    ("joint",13,True,DARK),(" problem — exactly what per-layer greedy methods cannot solve.",13,False,DARK)]])

# ============ SLIDE 12 : FINDING 3 — beats FLAT-LLM ============
s=base_slide()
txt(s,0.40,0.25,12.5,0.6,[[("Finding 3:  Amplification-Aware Allocation Beats FLAT-LLM",28,True,DARK)]])
rect(s,0.40,0.87,12.53,0.03,BLUE)
txt(s,12.50,7.10,0.70,0.30,[[("12",10,False,MUT)]],align=PP_ALIGN.RIGHT)
txt(s,0.40,1.02,12.5,0.4,[[("Same low-rank primitive, same total budget (",15,True,DARK),
    ("keep 50%",15,True,GREEN),("), Llama-2-7B WikiText perplexity — only the ",15,True,DARK),
    ("allocation",15,True,GREEN),(" differs.",15,True,DARK)]])
rows=[("No compression",6.87,GRAY),("Uniform",17.39,RGBColor(0x9E,0x9E,0x9E)),
      ("FLAT-LLM  (SOTA)",14.56,ORANGE),("Ours (amplification-joint)",13.96,GREEN)]
x0,y0,rowh,maxw=4.55,1.82,0.70,7.5; mx=max(v for _,v,_ in rows)
for i,(lab,v,col) in enumerate(rows):
    yy=y0+i*rowh; hot=("Ours" in lab or "FLAT" in lab)
    txt(s,0.55,yy+0.06,3.85,0.5,[[(lab,14,True,DARK if hot else GRAY)]],align=PP_ALIGN.RIGHT)
    w=maxw*v/mx
    rect(s,x0,yy,w,0.46,col,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s,x0+w+0.1,yy+0.04,1.4,0.4,[[(f"{v:.2f}",15,True,col)]])
txt(s,4.55,y0+4*rowh-0.02,7.5,0.3,[[("perplexity  (lower = better)",11,True,GRAY)]])
rect(s,0.40,5.05,6.15,0.92,LG); rect(s,0.40,5.05,0.16,0.92,GREEN)
txt(s,0.75,5.14,5.7,0.8,[[("−4.2%",26,True,GREEN)],[("vs FLAT-LLM  (faithful angular-BI baseline)",12,False,GRAY)]])
rect(s,6.75,5.05,6.15,0.92,LB); rect(s,6.75,5.05,0.16,0.92,BLUE)
txt(s,7.10,5.14,5.7,0.8,[[("−19.7%",26,True,BLUE)],[("vs Uniform allocation",12,False,GRAY)]])
rect(s,0.40,6.13,12.5,0.95,LY); rect(s,0.40,6.13,0.16,0.95,ORANGE)
txt(s,0.75,6.20,12.0,0.85,[[("Why we win — the signals disagree:  ",12,True,ORANGE),
    ("FLAT-LLM\u2019s angular BI protects boundary layers (0, 31, high residual rotation);",12,False,DARK)],
    [("amplification instead protects high-downstream-impact mid-early layers (2–5) and drains the damping tail.  ",12,False,DARK),
    ("Amp captures sensitivity BI misses.",12,True,DARK)]])

P.save("hypergraph_slides.pptx")
print("appended 3 slides; total now", len(list(P.slides)))
